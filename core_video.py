"""
core_video.py — 정확도 + 속도 + 안정성 업그레이드 버전

[정확도]
  1. Whisper `prompt` 도메인 힌트 주입
  2. 세그먼트 문장 단위 병합 (이중피동·어미 혼동 검출률↑)
  3. STT 오인식 가이드 + Few-shot + JSON Schema structured output
  4. OCR pHash 중복 제거

[속도]
  1. OCR 배치 병렬 호출 (ThreadPoolExecutor) — 2~5배 단축
  2. 프레임 추출 `cap.grab()` + 조건부 `retrieve()` 패턴
  3. ffmpeg 직접 호출 오디오 추출 (미설치 시 moviepy 폴백)
  4. JPEG 인코딩 품질/해상도 튜닝
  5. 음성·화면 파이프라인 병렬 실행 헬퍼

[안정성]
  1. 25MB 초과 오디오 자동 청크 분할 — 긴 영상 안전 처리
     + 청크별 병렬 STT + 원본 기준 타임스탬프 재조정
  2. Rate limit (429) 전용 처리 — Retry-After 헤더 존중, 더 긴 백오프
  3. STT 결과 캐싱 — 파일 해시 기반 (재실행 비용 0)
  4. 중간 실패 격리 — 한 배치 실패해도 나머지 유지 (이미 적용됨)
  5. 임시 파일 안전 관리 — context manager 제공

[정확도 v2] ★ 이번 업그레이드
  1. 스토리보드 내레이션 ground-truth 비교
     - STT 결과를 스토리보드 내레이션과 직접 비교 (단순 맞춤법 → 일치도 검증)
     - 내레이션 텍스트를 Whisper prompt에 주입해 STT 인식률 향상
  2. 히스토그램 + 구조적 유사도(SSIM) 기반 슬라이드 사전 매칭
     - GPT에 보내는 슬라이드 후보를 3~5장으로 좁혀 정밀도 향상 + 비용 절감
  3. 2-pass 검증 (화면)
     - 1차: 차이점 탐지 → 2차: 확신도 낮은 결과만 재검증 (false positive 감소)
  4. 음성-화면 교차 검증
     - 같은 시간대의 화면 텍스트와 음성을 함께 참조해 문맥 판단력 향상
  5. 프레임 추출 개선
     - 장면 전환 감지(scene change detection) 추가로 중요 프레임 놓침 방지
"""

import os
import re
import io
import json
import time
import base64
import hashlib
import shutil
import tempfile
import subprocess
from contextlib import contextmanager
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from typing import List, Optional, Callable, Dict, Any, Tuple

import cv2
import numpy as np
from PIL import Image as PILImage
from moviepy.editor import VideoFileClip
from openai import OpenAI
import win32com.client
import pythoncom


# ─────────────────────────────────────────────
# 유틸리티
# ─────────────────────────────────────────────

_FFMPEG_BIN = shutil.which("ffmpeg")


# ─────────────────────────────────────────────
# 이미지 유사도 — 슬라이드 사전 매칭 (정확도 v2)
# ─────────────────────────────────────────────

def _calc_histogram_similarity(img_b64_a: str, img_b64_b: str) -> float:
    """
    두 Base64 이미지의 히스토그램 상관도를 계산합니다.
    반환값: -1.0 ~ 1.0 (1.0이면 완전 일치)
    """
    try:
        buf_a = np.frombuffer(base64.b64decode(img_b64_a), np.uint8)
        buf_b = np.frombuffer(base64.b64decode(img_b64_b), np.uint8)
        img_a = cv2.imdecode(buf_a, cv2.IMREAD_COLOR)
        img_b = cv2.imdecode(buf_b, cv2.IMREAD_COLOR)
        if img_a is None or img_b is None:
            return -1.0

        # 동일 크기로 리사이즈
        target = (256, 256)
        img_a = cv2.resize(img_a, target)
        img_b = cv2.resize(img_b, target)

        # HSV 히스토그램 비교 (색상+밝기 모두 반영)
        hsv_a = cv2.cvtColor(img_a, cv2.COLOR_BGR2HSV)
        hsv_b = cv2.cvtColor(img_b, cv2.COLOR_BGR2HSV)

        h_bins, s_bins = 50, 60
        hist_a = cv2.calcHist([hsv_a], [0, 1], None, [h_bins, s_bins],
                              [0, 180, 0, 256])
        hist_b = cv2.calcHist([hsv_b], [0, 1], None, [h_bins, s_bins],
                              [0, 180, 0, 256])
        cv2.normalize(hist_a, hist_a)
        cv2.normalize(hist_b, hist_b)

        return cv2.compareHist(hist_a, hist_b, cv2.HISTCMP_CORREL)
    except Exception:
        return -1.0


def _calc_ssim_similarity(img_b64_a: str, img_b64_b: str) -> float:
    """
    두 Base64 이미지의 구조적 유사도(SSIM 근사)를 계산합니다.
    반환값: 0.0 ~ 1.0 (1.0이면 완전 일치)
    """
    try:
        buf_a = np.frombuffer(base64.b64decode(img_b64_a), np.uint8)
        buf_b = np.frombuffer(base64.b64decode(img_b64_b), np.uint8)
        img_a = cv2.imdecode(buf_a, cv2.IMREAD_GRAYSCALE)
        img_b = cv2.imdecode(buf_b, cv2.IMREAD_GRAYSCALE)
        if img_a is None or img_b is None:
            return 0.0

        target = (256, 256)
        img_a = cv2.resize(img_a, target).astype(np.float64)
        img_b = cv2.resize(img_b, target).astype(np.float64)

        # 간이 SSIM 계산 (scipy 없이)
        c1 = (0.01 * 255) ** 2
        c2 = (0.03 * 255) ** 2

        mu_a = cv2.GaussianBlur(img_a, (11, 11), 1.5)
        mu_b = cv2.GaussianBlur(img_b, (11, 11), 1.5)

        mu_a_sq = mu_a ** 2
        mu_b_sq = mu_b ** 2
        mu_ab = mu_a * mu_b

        sigma_a_sq = cv2.GaussianBlur(img_a ** 2, (11, 11), 1.5) - mu_a_sq
        sigma_b_sq = cv2.GaussianBlur(img_b ** 2, (11, 11), 1.5) - mu_b_sq
        sigma_ab = cv2.GaussianBlur(img_a * img_b, (11, 11), 1.5) - mu_ab

        ssim_map = ((2 * mu_ab + c1) * (2 * sigma_ab + c2)) / \
                   ((mu_a_sq + mu_b_sq + c1) * (sigma_a_sq + sigma_b_sq + c2))

        return float(np.mean(ssim_map))
    except Exception:
        return 0.0


# ─────────────────────────────────────────────
# ★ v5: 텍스트 콘텐츠 기반 슬라이드 매칭 (디자인 차이 무관)
# ─────────────────────────────────────────────
_KO_HANGUL_RE = re.compile(r'[가-힣A-Za-z0-9]+')

def _normalize_for_match(text: str) -> str:
    """매칭용 정규화: 공백·특수기호 제거, 소문자."""
    if not text:
        return ""
    # 한글·영문·숫자만 남김
    return "".join(_KO_HANGUL_RE.findall(text)).lower()


def _tokenize_for_match(text: str, min_len: int = 2) -> List[str]:
    """매칭용 토큰화: 한글/영문/숫자 토큰 분리, 짧은 토큰 제거, 중복 제거."""
    if not text:
        return []
    tokens = _KO_HANGUL_RE.findall(text)
    out: List[str] = []
    seen = set()
    for tok in tokens:
        t = tok.lower()
        if len(t) < min_len:
            continue
        if t in seen:
            continue
        seen.add(t)
        out.append(t)
    return out


def _text_match_score(frame_text: str, slide_texts: List[str]) -> Tuple[float, List[str]]:
    """
    영상 프레임 OCR 텍스트(transcription)와 슬라이드 본문 텍스트 리스트 간의
    텍스트 매칭 점수를 계산합니다.

    점수 구성 (0.0 ~ 1.0):
      · 슬라이드 본문 토큰 중 영상 프레임에 등장하는 비율 (recall)
      · 가중: 긴 토큰(고유명사·전문용어)일수록 점수 가산
      · 부분 일치(2/3 이상 글자 일치)도 인정

    Returns
    -------
    (score, matched_keywords)
    """
    if not frame_text or not slide_texts:
        return 0.0, []

    frame_norm = _normalize_for_match(frame_text)
    if not frame_norm:
        return 0.0, []

    # 슬라이드 토큰 모으기
    slide_tokens: List[str] = []
    for s in slide_texts:
        slide_tokens.extend(_tokenize_for_match(s, min_len=2))
    # 중복 제거, 길이 순 정렬 (긴 토큰 먼저 매칭 시도)
    seen = set()
    uniq_tokens: List[str] = []
    for t in slide_tokens:
        if t not in seen:
            seen.add(t)
            uniq_tokens.append(t)
    uniq_tokens.sort(key=len, reverse=True)

    if not uniq_tokens:
        return 0.0, []

    matched: List[str] = []
    weighted_hit = 0.0
    weighted_total = 0.0
    for tok in uniq_tokens:
        # 가중치: 긴 토큰일수록 정보량 높음 (제곱근)
        w = max(1.0, len(tok) ** 0.7)
        weighted_total += w
        # 1) 완전 포함
        if tok in frame_norm:
            weighted_hit += w
            matched.append(tok)
            continue
        # 2) 부분 포함 (3자 이상일 때만, 앞 2/3 또는 뒤 2/3 글자 매칭)
        if len(tok) >= 3:
            cut = max(2, (len(tok) * 2) // 3)
            if tok[:cut] in frame_norm or tok[-cut:] in frame_norm:
                weighted_hit += w * 0.6   # 부분 매칭은 60% 점수
                matched.append(tok + "*")

    score = weighted_hit / weighted_total if weighted_total > 0 else 0.0
    return score, matched


def text_based_match(
    frame_text: str,
    slide_metadata: List[dict],
    *,
    score_threshold: float = 0.10,
    margin: float = 0.05,
) -> Tuple[int, float, str, List[Tuple[int, float, List[str]]]]:
    """
    영상 프레임의 OCR 텍스트(transcription)를 슬라이드별 본문 텍스트와 비교하여
    가장 유사한 슬라이드의 1-based 번호를 반환합니다.

    ★ v5 → v6: 임계값을 0.30 → 0.20 으로 완화 (내레이션 포함으로 매칭 풀이 풍부해짐).

    Parameters
    ----------
    frame_text : 영상 프레임 OCR 결과 (모델의 transcription 필드)
    slide_metadata : build_slide_display_meta() 결과 (각 항목에 content_texts 포함)
    score_threshold : 이 점수 미만이면 매칭 없음으로 처리
    margin : 1등과 2등 점수 차가 이 미만이면 confidence 낮춤

    Returns
    -------
    (matched_slide_number, score, confidence, top_scores)
    - matched_slide_number : 1-based, 0이면 매칭 없음
    - score : 0.0 ~ 1.0
    - confidence : "exact" / "high" / "medium" / "low" / "none"
    - top_scores : 디버깅용 상위 후보 [(slide_no_1based, score, matched_keywords), ...] 최대 5개
    """
    if not frame_text or not slide_metadata:
        return 0, 0.0, "none", []

    scored: List[Tuple[int, float, List[str]]] = []
    for i, meta in enumerate(slide_metadata):
        cts = meta.get("content_texts") or []
        # 챕터·소주제도 매칭 텍스트에 포함 (영상에 보일 가능성 높음)
        extra = []
        if meta.get("subtopic"):
            extra.append(meta["subtopic"])
        if meta.get("title"):
            extra.append(meta["title"])
        if meta.get("chapter"):
            extra.append(meta["chapter"])
        score, matched = _text_match_score(frame_text, cts + extra)
        scored.append((i, score, matched))

    scored.sort(key=lambda x: x[1], reverse=True)
    best_idx, best_score, best_matched = scored[0]
    second_score = scored[1][1] if len(scored) > 1 else 0.0

    # 디버깅용 상위 5개 후보 (1-based 변환)
    top5 = [(idx + 1, sc, mk) for idx, sc, mk in scored[:5]]

    if best_score < score_threshold:
        return 0, best_score, "none", top5

    # confidence 등급 산정
    diff = best_score - second_score
    if best_score >= 0.60 and diff >= margin:
        conf = "exact"
    elif best_score >= 0.40 and diff >= margin:
        conf = "high"
    elif best_score >= 0.28:
        conf = "medium"
    else:
        conf = "low"

    return best_idx + 1, best_score, conf, top5


def find_best_matching_slides(
    frame_b64: str,
    storyboard_images: List[str],
    search_range: Optional[Tuple[int, int]] = None,
    top_k: int = 5,
) -> List[Tuple[int, float]]:
    """
    비디오 프레임과 가장 유사한 스토리보드 슬라이드를 찾습니다.

    히스토그램 상관도 + SSIM 가중 평균으로 순위 매김.
    GPT에 보내기 전 후보를 좁혀서 정확도와 비용을 동시에 개선합니다.

    Parameters
    ----------
    frame_b64 : str
        비디오 프레임 Base64
    storyboard_images : List[str]
        스토리보드 슬라이드 Base64 리스트
    search_range : Optional[Tuple[int, int]]
        탐색할 슬라이드 인덱스 범위 (start, end). None이면 전체.
    top_k : int
        반환할 상위 후보 수

    Returns
    -------
    List[(slide_index, combined_score)]
        유사도 점수가 높은 순으로 정렬된 (슬라이드 인덱스, 점수) 리스트
    """
    if not storyboard_images:
        return []

    start = search_range[0] if search_range else 0
    end = search_range[1] if search_range else len(storyboard_images)
    start = max(0, start)
    end = min(len(storyboard_images), end)

    scores: List[Tuple[int, float]] = []
    for i in range(start, end):
        hist_sim = _calc_histogram_similarity(frame_b64, storyboard_images[i])
        ssim_sim = _calc_ssim_similarity(frame_b64, storyboard_images[i])
        # 가중 평균: 히스토그램 40% + SSIM 60% (구조적 유사도에 더 비중)
        combined = hist_sim * 0.4 + ssim_sim * 0.6
        scores.append((i, combined))

    scores.sort(key=lambda x: x[1], reverse=True)
    return scores[:top_k]


def extract_audio(video_path, audio_path, sample_rate: int = 16000, mono: bool = True):
    """
    MP4 영상에서 오디오(.mp3)를 추출합니다.

    속도 우선:
    - ffmpeg가 설치돼 있으면 subprocess로 직접 호출 (moviepy 대비 3~5배 빠름)
    - 16kHz / mono 다운샘플로 업로드 크기 축소 + Whisper 정확도는 거의 영향 없음
    - ffmpeg 미설치 시 moviepy로 폴백

    Parameters
    ----------
    sample_rate : int
        오디오 샘플링 레이트 (Hz). Whisper는 16kHz로 내부 리샘플하므로 16000 권장.
    mono : bool
        모노 다운믹스 여부. 강의/대화 영상은 모노로 충분.
    """
    # 1) ffmpeg 경로가 있으면 직접 호출
    if _FFMPEG_BIN:
        try:
            cmd = [
                _FFMPEG_BIN,
                "-y",                     # 덮어쓰기
                "-loglevel", "error",
                "-i", video_path,
                "-vn",                    # 비디오 제거
                "-ar", str(sample_rate),
                "-ac", "1" if mono else "2",
                "-b:a", "64k",            # 64kbps 충분
                audio_path,
            ]
            subprocess.run(cmd, check=True, capture_output=True)
            return True
        except subprocess.CalledProcessError as e:
            stderr = (e.stderr or b"").decode("utf-8", errors="ignore")
            print(f"[ffmpeg 실패, moviepy로 폴백] {stderr.strip()[:200]}")
        except Exception as e:
            print(f"[ffmpeg 예외, moviepy로 폴백] {e}")

    # 2) 폴백: moviepy
    try:
        video = VideoFileClip(video_path)
        video.audio.write_audiofile(
            audio_path,
            fps=sample_rate,
            nbytes=2,
            codec="libmp3lame",
            bitrate="64k",
            ffmpeg_params=["-ac", "1" if mono else "2"],
            logger=None,
        )
        video.close()
        return True
    except Exception as e:
        print(f"오디오 추출 오류: {e}")
        return False


def format_timestamp(seconds):
    """초(float) 단위의 시간을 [HH:MM:SS] 또는 [MM:SS] 형태로 변환합니다."""
    hours = int(seconds // 3600)
    minutes = int((seconds % 3600) // 60)
    secs = int(seconds % 60)
    if hours > 0:
        return f"[{hours:02d}:{minutes:02d}:{secs:02d}]"
    return f"[{minutes:02d}:{secs:02d}]"


def call_with_retry(fn, retries=4, delay=3, max_delay: float = 60.0):
    """
    API 호출 실패 시 exponential backoff으로 재시도합니다.

    ★ 안정성 업그레이드:
    - RateLimitError(429) 감지 시 Retry-After 헤더를 존중 (있으면 그대로, 없으면 더 긴 대기)
    - APIConnectionError / Timeout 계열도 지수 백오프로 재시도
    - 최대 대기 시간(max_delay)으로 상한 둠

    retries : 총 시도 횟수 (기본 4)
    delay   : 초기 지연(초). 이후 2배씩 증가.
    """
    for attempt in range(retries):
        try:
            return fn()
        except Exception as e:
            err_name = type(e).__name__
            msg = str(e)
            is_rate_limit = "RateLimit" in err_name or "429" in msg

            # Retry-After 헤더 파싱 시도
            retry_after: Optional[float] = None
            try:
                resp = getattr(e, "response", None)
                if resp is not None:
                    hdr = getattr(resp, "headers", {}) or {}
                    ra = hdr.get("Retry-After") or hdr.get("retry-after")
                    if ra:
                        retry_after = float(ra)
            except Exception:
                pass

            if attempt < retries - 1:
                if is_rate_limit:
                    wait = retry_after if retry_after else min(max_delay, delay * (2 ** attempt) * 2)
                    print(f"⚠️ Rate limit ({err_name}) — {wait:.1f}초 대기 후 재시도 "
                          f"(시도 {attempt + 1}/{retries})")
                else:
                    wait = min(max_delay, delay * (2 ** attempt))
                    print(f"API 호출 오류 (시도 {attempt + 1}/{retries}) [{err_name}]: {msg[:200]}")
                time.sleep(wait)
            else:
                print(f"API 호출 최종 실패 [{err_name}]: {msg[:200]}")
    raise RuntimeError(f"API 호출이 {retries}회 모두 실패했습니다.")


# ─────────────────────────────────────────────
# 오디오 청크 분할 유틸 (25MB 초과 안전 처리)
# ─────────────────────────────────────────────

WHISPER_FILE_LIMIT_BYTES = 25 * 1024 * 1024   # 25 MB (Whisper API 제한)
# 실제로는 여유를 두고 20MB 기준으로 분할 판단
_CHUNK_SIZE_THRESHOLD = 20 * 1024 * 1024


def _get_audio_duration(audio_path: str) -> Optional[float]:
    """ffprobe로 오디오 길이(초)를 구합니다. 실패 시 None."""
    ffprobe = shutil.which("ffprobe")
    if not ffprobe:
        return None
    try:
        result = subprocess.run(
            [ffprobe, "-v", "error", "-show_entries", "format=duration",
             "-of", "default=noprint_wrappers=1:nokey=1", audio_path],
            capture_output=True, text=True, check=True, timeout=10,
        )
        return float(result.stdout.strip())
    except Exception:
        return None


def _split_audio_into_chunks(
    audio_path: str,
    chunk_duration: float,
    output_dir: str,
) -> List[Tuple[str, float]]:
    """
    ffmpeg로 오디오를 시간 기반으로 분할합니다.

    Returns
    -------
    List[(chunk_path, chunk_start_seconds)]
        각 청크 파일 경로와 원본 기준 시작 시각.
    """
    if not _FFMPEG_BIN:
        raise RuntimeError("오디오 청크 분할에는 ffmpeg가 필요합니다.")

    total_duration = _get_audio_duration(audio_path)
    if total_duration is None:
        # ffprobe 없어도 동작하게: 고정 개수로 나눌 수 없으므로 에러
        raise RuntimeError("ffprobe를 찾을 수 없어 오디오 길이를 확인할 수 없습니다.")

    chunks: List[Tuple[str, float]] = []
    idx = 0
    start = 0.0
    while start < total_duration:
        chunk_path = os.path.join(output_dir, f"chunk_{idx:03d}.mp3")
        # -ss는 입력 앞에 두면 빠른 seek, 뒤에 두면 정확한 seek.
        # 정확도를 위해 출력 옵션으로 사용.
        cmd = [
            _FFMPEG_BIN, "-y", "-loglevel", "error",
            "-i", audio_path,
            "-ss", f"{start:.3f}",
            "-t", f"{chunk_duration:.3f}",
            "-c", "copy",   # 재인코딩 없이 빠르게 복사
            chunk_path,
        ]
        try:
            subprocess.run(cmd, check=True, capture_output=True, timeout=60)
        except subprocess.CalledProcessError:
            # copy 실패 시 재인코딩으로 폴백
            cmd_reenc = [
                _FFMPEG_BIN, "-y", "-loglevel", "error",
                "-i", audio_path,
                "-ss", f"{start:.3f}",
                "-t", f"{chunk_duration:.3f}",
                "-ar", "16000", "-ac", "1", "-b:a", "64k",
                chunk_path,
            ]
            subprocess.run(cmd_reenc, check=True, capture_output=True, timeout=120)

        if os.path.exists(chunk_path) and os.path.getsize(chunk_path) > 0:
            chunks.append((chunk_path, start))
        idx += 1
        start += chunk_duration

    return chunks


def _plan_audio_chunks(audio_path: str, target_chunk_bytes: int = _CHUNK_SIZE_THRESHOLD) -> float:
    """
    파일 크기와 길이로부터 **청크당 적정 시간(초)** 을 계산합니다.
    """
    size = os.path.getsize(audio_path)
    duration = _get_audio_duration(audio_path) or 0.0
    if duration <= 0:
        # 안전 기본값: 10분
        return 600.0
    # 대략 "(목표_바이트 / 전체_바이트) * 전체_초"
    ratio = target_chunk_bytes / max(size, 1)
    chunk_sec = max(60.0, min(1500.0, duration * ratio * 0.95))   # 최소 1분, 최대 25분
    return chunk_sec


@contextmanager
def temp_workspace(prefix: str = "spellcheck_"):
    """
    임시 작업 디렉토리를 만들고, 블록 종료 시 자동 정리합니다.

    사용:
        with temp_workspace() as tmpdir:
            ...
    """
    tmpdir = tempfile.mkdtemp(prefix=prefix)
    try:
        yield tmpdir
    finally:
        try:
            shutil.rmtree(tmpdir, ignore_errors=True)
        except Exception:
            pass


# ─────────────────────────────────────────────
# 음성(STT) 처리
# ─────────────────────────────────────────────

@dataclass
class MergedSegment:
    """문장 단위로 병합된 세그먼트."""
    id: int
    start: float
    end: float
    text: str
    # 병합에 포함된 원본 Whisper 세그먼트 인덱스
    source_ids: List[int] = field(default_factory=list)


# 한국어 문장 종결 패턴 — 문장 병합 경계 판별용
_SENTENCE_END_PATTERN = re.compile(
    r"[.!?…]$|"
    r"(다|요|죠|네|까|군|구나|라|까요|어요|아요|예요|이에요|이다|입니다|습니다|니다)[.!?]?$"
)


def _is_sentence_end(text: str) -> bool:
    """문장이 끝났는지 추정 (한국어 종결어미 또는 문장부호)."""
    t = text.strip()
    if not t:
        return False
    return bool(_SENTENCE_END_PATTERN.search(t))


def merge_segments_into_sentences(
    segments,
    max_gap: float = 1.5,
    max_duration: float = 15.0,
    min_chars: int = 12,
):
    """
    Whisper의 짧은 세그먼트를 문장 단위로 병합합니다.

    병합 중단 조건:
    - 이전 세그먼트 끝과 현재 세그먼트 시작 간격(gap) > max_gap 초
    - 누적 길이(duration) >= max_duration 초
    - 이전 세그먼트가 문장 종결 형태이고 누적 길이가 min_chars 이상

    Parameters
    ----------
    segments : list
        Whisper verbose_json의 segments (obj 또는 dict)
    max_gap : float
        이 이상 침묵이면 강제로 문장 분리
    max_duration : float
        이 이상 길어지면 강제로 문장 분리 (너무 긴 입력 방지)
    min_chars : int
        최소 이 글자 수가 넘어야 문장 종료로 판정

    Returns
    -------
    List[MergedSegment]
    """
    def _get(seg, key):
        return getattr(seg, key, None) if not isinstance(seg, dict) else seg.get(key)

    merged: List[MergedSegment] = []
    buf_start: Optional[float] = None
    buf_end: Optional[float] = None
    buf_text = ""
    buf_sources: List[int] = []

    def _flush():
        nonlocal buf_start, buf_end, buf_text, buf_sources
        if buf_text.strip():
            merged.append(MergedSegment(
                id=len(merged),
                start=buf_start if buf_start is not None else 0.0,
                end=buf_end if buf_end is not None else 0.0,
                text=buf_text.strip(),
                source_ids=buf_sources.copy(),
            ))
        buf_start, buf_end, buf_text, buf_sources = None, None, "", []

    for i, seg in enumerate(segments):
        text = (_get(seg, "text") or "").strip()
        if not text:
            continue
        # Whisper가 뱉는 "." 만, "음…" 같은 무의미한 세그먼트 필터
        if len(text) <= 1 and not text.isalnum():
            continue

        start = float(_get(seg, "start") or 0.0)
        end = float(_get(seg, "end") or start)

        # 첫 버퍼 초기화
        if buf_start is None:
            buf_start, buf_end = start, end
            buf_text = text
            buf_sources = [i]
            continue

        gap = start - buf_end
        if gap > max_gap:
            _flush()
            buf_start, buf_end = start, end
            buf_text = text
            buf_sources = [i]
            continue

        # 누적 길이 체크
        if (end - buf_start) > max_duration:
            _flush()
            buf_start, buf_end = start, end
            buf_text = text
            buf_sources = [i]
            continue

        # 이어붙이기
        buf_text = f"{buf_text} {text}".strip()
        buf_end = end
        buf_sources.append(i)

        # 문장 종결 판정
        if _is_sentence_end(buf_text) and len(buf_text) >= min_chars:
            _flush()

    _flush()
    return merged


def _transcribe_single_file(client: OpenAI, audio_path: str, prompt: str):
    """단일 오디오 파일 전사. 내부용."""
    def _call():
        with open(audio_path, "rb") as audio_file:
            return client.audio.transcriptions.create(
                model="whisper-1",
                file=audio_file,
                language="ko",
                response_format="verbose_json",
                timestamp_granularities=["segment", "word"],
                prompt=prompt,
                temperature=0.0,
            )
    return call_with_retry(_call)


def _shift_segments(segments, offset: float):
    """세그먼트의 start/end에 offset(초)를 더해 원본 타임라인으로 복원."""
    shifted = []
    for seg in segments:
        is_dict = isinstance(seg, dict)
        start = (seg.get("start") if is_dict else getattr(seg, "start", 0)) or 0
        end   = (seg.get("end")   if is_dict else getattr(seg, "end",   0)) or 0
        text  = (seg.get("text")  if is_dict else getattr(seg, "text", "")) or ""
        shifted.append(_ShiftedSeg(
            start=float(start) + offset,
            end=float(end) + offset,
            text=text,
        ))
    return shifted


@dataclass
class _ShiftedSeg:
    """청크 병합용 경량 세그먼트 (merge_segments_into_sentences와 호환)."""
    start: float
    end: float
    text: str


def _extract_narration_keywords(narration_texts: List[str], max_chars: int = 200) -> str:
    """
    스토리보드 내레이션 텍스트에서 Whisper prompt용 핵심 키워드를 추출합니다.

    전략:
    - 3글자 이상의 고유명사/전문용어 후보를 빈도순으로 선별
    - 일반적인 조사·어미를 제거하고 명사 위주로 추출
    - 총 max_chars 이내로 제한 (Whisper prompt 토큰 제한 대응)
    """
    if not narration_texts:
        return ""

    all_text = " ".join(t for t in narration_texts if t)

    # 한글 2글자 이상 단어 추출 (조사 제거를 위해 어절 단위 처리)
    # 간단한 형태소 근사: 어절에서 흔한 조사/어미를 제거
    _SUFFIX_PATTERN = re.compile(
        r"(은|는|이|가|을|를|에|에서|로|으로|와|과|의|도|만|까지|부터|처럼|보다|라고|라는|"
        r"하는|하고|해서|하여|하면|했다|합니다|입니다|있다|없다|됩니다|입니까|인데|이며|이고)$"
    )

    word_freq: Dict[str, int] = {}
    for word in re.findall(r"[가-힣a-zA-Z0-9]+", all_text):
        if len(word) < 2:
            continue
        # 조사/어미 제거
        cleaned = _SUFFIX_PATTERN.sub("", word)
        if len(cleaned) >= 2:
            word_freq[cleaned] = word_freq.get(cleaned, 0) + 1

    # 빈도 2 이상이거나 3글자 이상인 단어 우선 (일반적인 단어는 제외)
    _COMMON_WORDS = {
        "그리고", "하지만", "그래서", "이것", "저것", "우리", "여기", "거기",
        "이런", "저런", "그런", "때문", "다음", "먼저", "지금", "오늘",
        "내용", "부분", "경우", "사실", "정도", "이상", "이하", "통해",
    }
    candidates = [
        (w, c) for w, c in word_freq.items()
        if w not in _COMMON_WORDS and (c >= 2 or len(w) >= 3)
    ]
    candidates.sort(key=lambda x: (-x[1], -len(x[0])))

    result = []
    total_len = 0
    for word, _ in candidates:
        if total_len + len(word) + 2 > max_chars:
            break
        result.append(word)
        total_len += len(word) + 2  # ", " separator

    return ", ".join(result)


def transcribe_audio(
    audio_path,
    api_key,
    model: str = "whisper-1",
    domain_hint: str = "",
    max_chunk_bytes: int = _CHUNK_SIZE_THRESHOLD,
    max_workers: int = 3,
    narration_texts: Optional[List[str]] = None,
):
    """
    오디오 → 텍스트 변환.

    ★ 안정성 업그레이드:
    - 파일이 Whisper 25MB 제한을 초과하면 자동으로 ffmpeg로 청크 분할 후 병렬 전송
    - 각 청크 결과의 타임스탬프를 원본 기준으로 재조정해 이어 붙임
    - 한 청크가 실패해도 나머지 청크 결과는 보존 (부분 복구)

    ★ 정확도 v2:
    - narration_texts: 스토리보드 내레이션 텍스트 리스트를 Whisper prompt에 주입하여
      고유명사, 전문용어 등의 인식 정확도를 높입니다.
      Whisper prompt는 244토큰 제한이므로 핵심 문구를 추출하여 사용합니다.

    Parameters
    ----------
    max_chunk_bytes : int
        청크 분할 임계값. 기본 20MB (25MB 한도에 여유 둠).
    max_workers : int
        청크 병렬 전사 워커 수. 권장 2~4.
    narration_texts : Optional[List[str]]
        스토리보드 내레이션 텍스트 리스트. Whisper prompt 강화용.

    Returns
    -------
    segments : Whisper verbose_json segments (원본 타임라인으로 복원됨)
    """
    client = OpenAI(api_key=api_key)

    if model != "whisper-1":
        print(f"[안내] {model}은 세그먼트 타임스탬프를 지원하지 않아 "
              f"whisper-1로 자동 폴백합니다. (정확도 보완은 프롬프트 힌트로 유지)")

    prompt = (
        "다음은 한국어 강의 또는 발표 영상의 음성입니다. "
        "정확한 한국어 맞춤법과 고유명사 표기를 사용해 전사하세요."
    )
    if domain_hint.strip():
        prompt += f" 주요 용어: {domain_hint.strip()}"

    # ★ 정확도 v2: 내레이션 텍스트에서 핵심 키워드를 추출해 prompt에 추가
    # Whisper prompt는 ~244토큰 제한이므로 핵심 용어만 선별
    if narration_texts:
        key_phrases = _extract_narration_keywords(narration_texts)
        if key_phrases:
            prompt += f" 스토리보드 핵심 용어: {key_phrases}"

    # ── 파일 크기 확인 ─────────────────────────────
    try:
        file_size = os.path.getsize(audio_path)
    except OSError as e:
        raise RuntimeError(f"오디오 파일을 읽을 수 없습니다: {e}")

    # 25MB 이하면 단일 호출
    if file_size <= max_chunk_bytes:
        transcription = _transcribe_single_file(client, audio_path, prompt)
        return transcription.segments

    # ── 25MB 초과: 청크 분할 ───────────────────────
    size_mb = file_size / (1024 * 1024)
    print(f"📦 오디오 크기 {size_mb:.1f}MB → Whisper 25MB 한도 초과, 청크 분할 시작")

    if not _FFMPEG_BIN:
        raise RuntimeError(
            f"오디오가 {size_mb:.1f}MB로 Whisper 25MB 한도를 초과합니다. "
            f"청크 분할을 위해 ffmpeg 설치가 필요합니다."
        )

    chunk_seconds = _plan_audio_chunks(audio_path, max_chunk_bytes)
    print(f"   → 청크 길이: {chunk_seconds:.0f}초 (약 {chunk_seconds/60:.1f}분)")

    all_segments: List[Any] = []
    with temp_workspace(prefix="stt_chunks_") as tmpdir:
        try:
            chunks = _split_audio_into_chunks(audio_path, chunk_seconds, tmpdir)
        except Exception as e:
            raise RuntimeError(f"오디오 청크 분할 실패: {e}")

        if not chunks:
            raise RuntimeError("오디오 청크 분할 결과가 비어 있습니다.")

        print(f"   → 총 {len(chunks)}개 청크 생성, 병렬 전사 시작 (workers={max_workers})")

        def _process_chunk(chunk_path: str, offset: float, idx: int):
            try:
                tr = _transcribe_single_file(client, chunk_path, prompt)
                shifted = _shift_segments(tr.segments, offset)
                return idx, shifted, None
            except Exception as e:
                return idx, [], str(e)

        indexed_results: Dict[int, List[Any]] = {}
        errors: List[str] = []

        if max_workers <= 1 or len(chunks) == 1:
            for i, (cp, off) in enumerate(chunks):
                idx, segs, err = _process_chunk(cp, off, i)
                indexed_results[idx] = segs
                if err:
                    errors.append(f"청크 {idx+1}: {err}")
        else:
            with ThreadPoolExecutor(max_workers=max_workers) as ex:
                futures = {
                    ex.submit(_process_chunk, cp, off, i): i
                    for i, (cp, off) in enumerate(chunks)
                }
                for fut in as_completed(futures):
                    idx, segs, err = fut.result()
                    indexed_results[idx] = segs
                    if err:
                        errors.append(f"청크 {idx+1}: {err}")

        # 시간 순서대로 이어붙임
        for i in sorted(indexed_results.keys()):
            all_segments.extend(indexed_results[i])

    if errors:
        print(f"⚠️ 일부 청크 실패 ({len(errors)}개): {errors[0][:200]}")
        if not all_segments:
            raise RuntimeError(f"모든 청크 전사 실패: {errors[0]}")

    print(f"   ✅ 청크 전사 완료: {len(all_segments)}개 세그먼트")
    return all_segments


# ─────────────────────────────────────────────
# STT 결과 캐싱 (파일 해시 기반)
# ─────────────────────────────────────────────

def _file_content_hash(path: str, chunk_size: int = 1024 * 1024) -> str:
    """파일 SHA-256 해시. 큰 파일도 스트리밍으로 처리."""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        while True:
            chunk = f.read(chunk_size)
            if not chunk:
                break
            h.update(chunk)
    return h.hexdigest()


def compute_stt_cache_key(video_path: str, domain_hint: str = "") -> str:
    """
    STT 결과 캐싱용 키를 만듭니다.
    (파일 내용 해시 + 도메인 힌트)
    """
    file_hash = _file_content_hash(video_path)
    hint_hash = hashlib.md5(domain_hint.encode("utf-8")).hexdigest()[:8]
    return f"{file_hash[:16]}_{hint_hash}"


def save_stt_cache(cache_dir: str, cache_key: str, segments) -> bool:
    """전사 결과를 JSON으로 저장 (세그먼트는 start/end/text만 보존)."""
    try:
        os.makedirs(cache_dir, exist_ok=True)
        path = os.path.join(cache_dir, f"stt_{cache_key}.json")
        simplified = []
        for s in segments:
            is_dict = isinstance(s, dict)
            simplified.append({
                "start": float((s.get("start") if is_dict else getattr(s, "start", 0)) or 0),
                "end":   float((s.get("end")   if is_dict else getattr(s, "end",   0)) or 0),
                "text":  (s.get("text") if is_dict else getattr(s, "text", "")) or "",
            })
        with open(path, "w", encoding="utf-8") as f:
            json.dump(simplified, f, ensure_ascii=False)
        return True
    except Exception as e:
        print(f"STT 캐시 저장 실패: {e}")
        return False


def load_stt_cache(cache_dir: str, cache_key: str):
    """캐시된 전사 결과를 로드. 없으면 None."""
    path = os.path.join(cache_dir, f"stt_{cache_key}.json")
    if not os.path.exists(path):
        return None
    try:
        with open(path, "r", encoding="utf-8") as f:
            raw = json.load(f)
        return [_ShiftedSeg(start=x["start"], end=x["end"], text=x["text"]) for x in raw]
    except Exception as e:
        print(f"STT 캐시 로드 실패: {e}")
        return None


# ─────────────────────────────────────────────
# 맞춤법 검사 — 공통 규칙 & 프롬프트
# ─────────────────────────────────────────────

_KO_SPELL_RULES = """
주요 검사 항목 (특히 집중):
1. 사이시옷: 숫자·뒷말 등 합성어 표기 (예: 나뭇잎, 숫자)
2. 된소리·거센소리 혼동: '깨끗이' vs '깨끗히', '않다' vs '안다'
3. 어미 혼동: '-이에요'/'-이어요', '-데'/'-대', '-든지'/'-던지'
4. 조사·의존명사 띄어쓰기: '것', '수', '때', '만큼', '뿐' 등
5. 피동·사동 표현 오용: '되어지다', '만들어지다' 등 이중피동
6. 외래어 표기: '컨텐츠→콘텐츠', '메세지→메시지' 등
7. 공식 명칭·고유어 오기: 틀린 한자어·혼용 표기
8. 불필요한 중복 표현: '미리 예방', '과반수 이상' 등
9. 단위·숫자 표기: '1 개월' 등 띄어쓰기
10. 어순 및 비문(문장 구조 오류)
"""

# STT 특성을 반영한 가이드 — false positive 감소
_STT_CAVEAT = """
【STT 변환 특성 주의사항】
이 텍스트는 음성인식(STT)으로 변환되어, 아래 경우에는 원문을 유지하세요:
- 발화자의 말버릇·구어체(예: "어", "음", "그니까")는 맞춤법 오류가 아님
- 외래어가 한글로 음차된 경우 문맥상 맞으면 유지
- 동음이의어 선택이 애매하면 원문 유지
- STT 경계가 어색해도 의미가 통하면 유지

반대로 **확실한 오류**는 반드시 교정:
- 명백한 표준어 규정 위반 (예: "됬다" → "됐다")
- 외래어 표기법 위반 (예: "컨텐츠" → "콘텐츠")
- 이중피동 (예: "되어지다" → "되다")
"""

_SPELL_FEWSHOT = """
【교정 예시】
원문: "영상 컨텐츠가 만들어지면 됩니다."
교정: "영상 <red>콘텐츠</red>가 <red>만들어지면</red> 됩니다."
      → 콘텐츠(외래어 표기), 만들어지면(이중피동 주의 — 문맥에 따라 '만들면')

원문: "해당 내용은 다음주에 발표할께요."
교정: "해당 내용은 <red>다음 주</red>에 <red>발표할게요</red>."
      → 다음 주(의존명사 띄어쓰기), 발표할게요(어미 '-ㄹ게')

원문: "과반수 이상이 찬성했어요."
교정: "<red>과반수가</red> 찬성했어요."
      → '과반수' 자체가 절반 초과 의미 → '이상' 중복 표현
"""

_SPELL_SYSTEM_PROMPT = f"""당신은 대한민국 방송 표준어 규정과 한글 맞춤법을 완벽히 숙지한 수석 교열 전문가입니다.
아래 텍스트는 동영상 음성을 STT로 변환한 결과입니다. (각 줄에 ID, 시간, 텍스트 포함)
일부 줄은 [CONTEXT] 표시가 되어 있으며, 이는 앞뒤 문맥 참고용으로만 사용하고 **교정 대상이 아닙니다**.

{_KO_SPELL_RULES}
{_STT_CAVEAT}
{_SPELL_FEWSHOT}

【작업 지시】
- [CONTEXT] 줄은 참고만 하고 결과에 포함하지 마세요.
- 전체 문맥을 파악한 뒤 교정하세요.
- 오류가 확실할 때만 교정. 애매하면 원문 유지.
- 수정된 단어·어절만 <red>단어</red>로 감싸세요.
- 원문 문장 전체를 반환하되 수정 부분만 태그 처리.
- 오류 없는 세그먼트는 결과 배열에서 제외.
"""

# Structured Output용 JSON Schema
_SPELL_SCHEMA = {
    "name": "spell_corrections",
    "strict": True,
    "schema": {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "corrections": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "id":        {"type": "integer"},
                        "original":  {"type": "string"},
                        "corrected": {"type": "string"},
                        "reason":    {"type": "string"},
                    },
                    "required": ["id", "original", "corrected", "reason"],
                },
            }
        },
        "required": ["corrections"],
    },
}


# ─────────────────────────────────────────────
# 음성 맞춤법 검사 (세그먼트 병합 + 슬라이딩 윈도우)
# ─────────────────────────────────────────────

def spell_check_segments(
    segments,
    api_key,
    context_window: int = 2,
    model: str = "gpt-5.4",
    use_sentence_merge: bool = True,
    batch_size: int = 40,
    max_workers: int = 3,
):
    """
    세그먼트를 문장 단위로 병합한 뒤 맞춤법 검사합니다.

    Parameters
    ----------
    context_window : int
        각 배치 앞/뒤로 덧붙일 참고 세그먼트 수 (교정 대상 아님)
    use_sentence_merge : bool
        True면 Whisper 세그먼트를 문장 단위로 병합 (정확도↑)
    batch_size : int
        한 번에 모델에 보낼 교정 대상 세그먼트 수
    max_workers : int
        ★ 속도 업그레이드: 배치 병렬 호출 개수. OpenAI rate limit 고려 권장 2~4.
    """
    client = OpenAI(api_key=api_key)

    # 1) 문장 단위 병합
    if use_sentence_merge:
        merged = merge_segments_into_sentences(segments)
    else:
        merged = [
            MergedSegment(
                id=i,
                start=float(getattr(s, "start", s.get("start", 0)) if hasattr(s, "start") or isinstance(s, dict) else 0),
                end=float(getattr(s, "end", s.get("end", 0)) if hasattr(s, "end") or isinstance(s, dict) else 0),
                text=(getattr(s, "text", s.get("text", "")) if hasattr(s, "text") or isinstance(s, dict) else "").strip(),
                source_ids=[i],
            )
            for i, s in enumerate(segments)
            if (getattr(s, "text", s.get("text", "")) if hasattr(s, "text") or isinstance(s, dict) else "").strip()
        ]

    if not merged:
        return []

    # 2) 배치 생성 (컨텍스트 포함 텍스트 미리 빌드)
    batches: List[Tuple[int, int, str]] = []  # (batch_start, batch_end, batch_text)
    for batch_start in range(0, len(merged), batch_size):
        batch_end = min(batch_start + batch_size, len(merged))
        ctx_start = max(0, batch_start - context_window)
        ctx_end = min(len(merged), batch_end + context_window)

        lines = []
        for i in range(ctx_start, ctx_end):
            seg = merged[i]
            tag = "[CONTEXT] " if (i < batch_start or i >= batch_end) else ""
            lines.append(f"{tag}{format_timestamp(seg.start)} (ID:{seg.id}): {seg.text}")
        batches.append((batch_start, batch_end, "\n".join(lines)))

    def _process_batch(batch_idx: int, bs: int, be: int, text: str) -> List[dict]:
        """단일 배치 처리 — 병렬 워커용."""
        def _call():
            return client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": _SPELL_SYSTEM_PROMPT},
                    {"role": "user", "content": text},
                ],
                temperature=0.0,
                response_format={"type": "json_schema", "json_schema": _SPELL_SCHEMA},
            )

        batch_results = []
        try:
            response = call_with_retry(_call)
            parsed = json.loads(response.choices[0].message.content)
            corrections = parsed.get("corrections", [])

            for item in corrections:
                idx = item.get("id")
                if idx is None:
                    continue
                seg = next((m for m in merged if m.id == idx), None)
                if seg is None:
                    continue
                orig = (item.get("original") or "").strip()
                corr = (item.get("corrected") or "").strip()
                if not (orig and corr and orig != corr):
                    continue
                batch_results.append({
                    "구분": "음성 대본",
                    "시간": format_timestamp(seg.start),
                    "수정 전": _red_to_html(orig),
                    "수정 후": _red_to_html(corr),
                    "교정 사유": item.get("reason") or "",
                })
        except Exception as e:
            print(f"음성 맞춤법 검사 오류 (배치 {batch_idx + 1}): {e}")
            try:
                batch_results.extend(_fallback_text_mode_spell_check(
                    client, model, merged, bs, be, context_window
                ))
            except Exception as e2:
                print(f"  폴백도 실패: {e2}")
        return batch_results

    # 3) 병렬 실행
    results: List[dict] = []
    if max_workers <= 1 or len(batches) == 1:
        for i, (bs, be, text) in enumerate(batches):
            results.extend(_process_batch(i, bs, be, text))
    else:
        with ThreadPoolExecutor(max_workers=max_workers) as ex:
            futures = {
                ex.submit(_process_batch, i, bs, be, text): i
                for i, (bs, be, text) in enumerate(batches)
            }
            # 시간순 정렬을 위해 인덱스별로 받은 뒤 합치기
            indexed: Dict[int, List[dict]] = {}
            for fut in as_completed(futures):
                batch_idx = futures[fut]
                try:
                    indexed[batch_idx] = fut.result()
                except Exception as e:
                    print(f"병렬 배치 {batch_idx + 1} 실패: {e}")
                    indexed[batch_idx] = []
            for i in sorted(indexed.keys()):
                results.extend(indexed[i])

    return results


def _fallback_text_mode_spell_check(client, model, merged, batch_start, batch_end, context_window):
    """json_schema를 지원하지 않는 모델용 폴백 (기존 방식)."""
    ctx_start = max(0, batch_start - context_window)
    ctx_end = min(len(merged), batch_end + context_window)

    lines = []
    for i in range(ctx_start, ctx_end):
        seg = merged[i]
        tag = "[CONTEXT] " if (i < batch_start or i >= batch_end) else ""
        lines.append(f"{tag}{format_timestamp(seg.start)} (ID:{seg.id}): {seg.text}")

    instruction = (
        _SPELL_SYSTEM_PROMPT
        + "\n\n응답은 반드시 아래 순수 JSON 배열 형식으로만(백틱·설명 없이):\n"
        '[{"id":0,"original":"…","corrected":"…","reason":"…"}]\n오류 없으면 []'
    )

    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": instruction},
            {"role": "user", "content": "\n".join(lines)},
        ],
        temperature=0.0,
    )
    content = _strip_json_fences(response.choices[0].message.content.strip())
    corrections = json.loads(content)
    if isinstance(corrections, dict):
        corrections = corrections.get("corrections", next(iter(corrections.values()), []))

    out = []
    for item in corrections:
        idx = item.get("id")
        if idx is None:
            continue
        seg = next((m for m in merged if m.id == idx), None)
        if seg is None:
            continue
        orig = (item.get("original") or "").strip()
        corr = (item.get("corrected") or "").strip()
        if orig and corr and orig != corr:
            out.append({
                "구분": "음성 대본",
                "시간": format_timestamp(seg.start),
                "수정 전": _red_to_html(orig),
                "수정 후": _red_to_html(corr),
                "교정 사유": item.get("reason", ""),
            })
    return out


# ─────────────────────────────────────────────
# ★ 정확도 v2: 스토리보드 내레이션 대조 검증
# ─────────────────────────────────────────────

_NARRATION_COMPARE_PROMPT = f"""당신은 영상 음성과 스토리보드 원본 내레이션의 일치 여부를 정밀하게 검수하는 전문 교열자입니다.
대한민국 방송 표준어 규정과 한글 맞춤법을 완벽히 숙지하고 있습니다.

【입력 형식】
- [STT] 줄: 영상 음성을 자동 전사한 텍스트 (타임스탬프 + ID 포함)
- [REF] 줄: 해당 시점에 대응하는 스토리보드 원본 내레이션 (ground truth)
- [CONTEXT] 줄: 참고용 앞뒤 문맥 (교정 대상 아님)

【작업 지시】
1. 각 STT 텍스트를 대응하는 REF 내레이션과 면밀히 비교하세요.
2. 다음 유형의 차이를 찾으세요:
   a) **누락(missing)**: REF에 있는 핵심 내용이 STT에서 빠진 경우
   b) **오인식(mismatch)**: STT가 단어를 잘못 인식한 경우 (예: "연구" → "영구")
   c) **추가(addition)**: STT에 있지만 REF에 없는 의미상 중요한 내용
   d) **맞춤법(spelling)**: REF 기준으로 맞춤법이 다른 경우
3. **주의 — 다음은 오류가 아닙니다**:
   - 어순 변경, 구어체 변형(예: "합니다"→"해요")
   - 자연스러운 말줄임, 추임새("음", "어", "그")
   - 동의어 치환(의미가 동일한 경우)
4. REF가 없는 STT 줄은 아래 맞춤법 규칙에 따라 일반 맞춤법 검사만 수행하세요.
5. confidence: "high"(확실), "medium"(거의 확실), "low"(가능성 있음) 중 택1

{_KO_SPELL_RULES}
{_STT_CAVEAT}
{_SPELL_FEWSHOT}

【출력 규칙】
- original: STT 텍스트에 나온 그대로 (문장 전체, 수정 부분만 <red>태그</red>)
- corrected: REF 기준의 올바른 표현, 수정 부분만 <red>태그</red>
- reason: "[내레이션 불일치] ~" 또는 "[맞춤법] ~" 형식으로 명확히 분류
- type: "mismatch" | "missing" | "addition" | "spelling" 중 택1
- 차이 없는 세그먼트는 결과에서 제외
"""

_NARRATION_COMPARE_SCHEMA = {
    "name": "narration_corrections",
    "strict": True,
    "schema": {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "corrections": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "id":         {"type": "integer"},
                        "original":   {"type": "string"},
                        "corrected":  {"type": "string"},
                        "reason":     {"type": "string"},
                        "type":       {"type": "string"},
                        "confidence": {"type": "string"},
                    },
                    "required": ["id", "original", "corrected", "reason", "type", "confidence"],
                },
            }
        },
        "required": ["corrections"],
    },
}


def _align_stt_to_narration(
    merged_segments: List[MergedSegment],
    narration_texts: List[str],
) -> List[Tuple[int, Optional[str]]]:
    """
    STT 세그먼트를 내레이션 텍스트와 텍스트 유사도 기반으로 매칭합니다.

    내레이션이 슬라이드 단위로 주어지므로 여러 STT 세그먼트가
    하나의 내레이션에 매칭될 수 있습니다.

    Returns
    -------
    List[(segment_id, matched_narration_text or None)]
    """
    if not narration_texts:
        return [(seg.id, None) for seg in merged_segments]

    # 자카드 유사도
    def _jaccard(a: str, b: str) -> float:
        set_a = set(re.findall(r"[가-힣a-zA-Z0-9]+", a))
        set_b = set(re.findall(r"[가-힣a-zA-Z0-9]+", b))
        if not set_a or not set_b:
            return 0.0
        return len(set_a & set_b) / len(set_a | set_b)

    result = []
    narr_idx = 0

    for seg in merged_segments:
        if narr_idx >= len(narration_texts):
            result.append((seg.id, None))
            continue

        best_score = -1.0
        best_narr = None
        best_ni = narr_idx
        for offset in range(-1, 4):
            ni = narr_idx + offset
            if 0 <= ni < len(narration_texts) and narration_texts[ni]:
                score = _jaccard(seg.text, narration_texts[ni])
                if score > best_score:
                    best_score = score
                    best_narr = narration_texts[ni]
                    best_ni = ni

        if best_score >= 0.1 and best_narr:
            result.append((seg.id, best_narr))
            if best_ni > narr_idx:
                narr_idx = best_ni
        else:
            result.append((seg.id, None))

    return result


def narration_compare_segments(
    segments,
    api_key: str,
    narration_texts: List[str],
    context_window: int = 2,
    model: str = "gpt-5.4",
    use_sentence_merge: bool = True,
    batch_size: int = 30,
    max_workers: int = 3,
    confidence_threshold: str = "medium",
) -> List[dict]:
    """
    ★ 정확도 v2 핵심 기능
    STT 결과를 스토리보드 내레이션과 대조하여 불일치를 찾습니다.

    기존 spell_check_segments는 일반 맞춤법만 검사했지만,
    이 함수는 내레이션 원문(ground truth)과 비교하여:
    - 오인식 단어 탐지
    - 누락/추가 탐지
    - 맞춤법 오류도 함께 검사

    Parameters
    ----------
    narration_texts : List[str]
        스토리보드 슬라이드별 내레이션 텍스트
    confidence_threshold : str
        "high" | "medium" | "low". 이 수준 이상만 결과에 포함.
    """
    client = OpenAI(api_key=api_key)

    # 1) 문장 단위 병합
    if use_sentence_merge:
        merged = merge_segments_into_sentences(segments)
    else:
        merged = [
            MergedSegment(
                id=i,
                start=float(getattr(s, "start", s.get("start", 0)) if hasattr(s, "start") or isinstance(s, dict) else 0),
                end=float(getattr(s, "end", s.get("end", 0)) if hasattr(s, "end") or isinstance(s, dict) else 0),
                text=(getattr(s, "text", s.get("text", "")) if hasattr(s, "text") or isinstance(s, dict) else "").strip(),
                source_ids=[i],
            )
            for i, s in enumerate(segments)
            if (getattr(s, "text", s.get("text", "")) if hasattr(s, "text") or isinstance(s, dict) else "").strip()
        ]

    if not merged:
        return []

    # 2) STT ↔ 내레이션 정렬
    alignment = _align_stt_to_narration(merged, narration_texts)
    narr_map = {seg_id: narr for seg_id, narr in alignment}

    # 3) 배치 생성 (컨텍스트 + 내레이션 참조 포함)
    _conf_levels = {"high": 3, "medium": 2, "low": 1}
    threshold_level = _conf_levels.get(confidence_threshold, 2)

    batches: List[Tuple[int, int, str]] = []
    for batch_start in range(0, len(merged), batch_size):
        batch_end = min(batch_start + batch_size, len(merged))
        ctx_start = max(0, batch_start - context_window)
        ctx_end = min(len(merged), batch_end + context_window)

        lines = []
        for i in range(ctx_start, ctx_end):
            seg = merged[i]
            tag = "[CONTEXT] " if (i < batch_start or i >= batch_end) else ""
            lines.append(f"{tag}[STT] {format_timestamp(seg.start)} (ID:{seg.id}): {seg.text}")
            narr = narr_map.get(seg.id)
            if narr and not (i < batch_start or i >= batch_end):
                lines.append(f"  [REF] {narr}")

        batches.append((batch_start, batch_end, "\n".join(lines)))

    def _process_batch(batch_idx: int, bs: int, be: int, text: str) -> List[dict]:
        def _call():
            return client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": _NARRATION_COMPARE_PROMPT},
                    {"role": "user", "content": text},
                ],
                temperature=0.0,
                response_format={"type": "json_schema", "json_schema": _NARRATION_COMPARE_SCHEMA},
            )

        batch_results = []
        try:
            response = call_with_retry(_call)
            parsed = json.loads(response.choices[0].message.content)
            corrections = parsed.get("corrections", [])

            for item in corrections:
                idx = item.get("id")
                if idx is None:
                    continue
                seg = next((m for m in merged if m.id == idx), None)
                if seg is None:
                    continue
                orig = (item.get("original") or "").strip()
                corr = (item.get("corrected") or "").strip()
                if not (orig and corr and orig != corr):
                    continue

                conf = item.get("confidence", "medium")
                conf_level = _conf_levels.get(conf, 2)
                if conf_level < threshold_level:
                    continue

                err_type = item.get("type", "spelling")
                reason_prefix = {
                    "mismatch": "[내레이션 불일치]",
                    "missing": "[내레이션 누락]",
                    "addition": "[내레이션 추가]",
                    "spelling": "[맞춤법]",
                }.get(err_type, "")

                reason = item.get("reason") or ""
                if reason_prefix and not reason.startswith(reason_prefix):
                    reason = f"{reason_prefix} {reason}"

                batch_results.append({
                    "구분": "음성 대본",
                    "시간": format_timestamp(seg.start),
                    "수정 전": _red_to_html(orig),
                    "수정 후": _red_to_html(corr),
                    "교정 사유": reason,
                    "_confidence": conf,
                    "_type": err_type,
                })
        except Exception as e:
            print(f"내레이션 비교 오류 (배치 {batch_idx + 1}): {e}")
            try:
                batch_results.extend(_fallback_text_mode_spell_check(
                    client, model, merged, bs, be, context_window
                ))
            except Exception as e2:
                print(f"  폴백도 실패: {e2}")
        return batch_results

    # 4) 병렬 실행
    results: List[dict] = []
    if max_workers <= 1 or len(batches) == 1:
        for i, (bs, be, text) in enumerate(batches):
            results.extend(_process_batch(i, bs, be, text))
    else:
        with ThreadPoolExecutor(max_workers=max_workers) as ex:
            futures = {
                ex.submit(_process_batch, i, bs, be, text): i
                for i, (bs, be, text) in enumerate(batches)
            }
            indexed: Dict[int, List[dict]] = {}
            for fut in as_completed(futures):
                batch_idx = futures[fut]
                try:
                    indexed[batch_idx] = fut.result()
                except Exception as e:
                    print(f"병렬 배치 {batch_idx + 1} 실패: {e}")
                    indexed[batch_idx] = []
            for i in sorted(indexed.keys()):
                results.extend(indexed[i])

    return results


# ─────────────────────────────────────────────
# 화면 프레임 처리
# ─────────────────────────────────────────────

# 전역 품질 설정 (속도 모드에서 app 쪽에서 오버라이드 가능)
_JPEG_QUALITY = 90           # 기존 95 → 90으로 낮춰도 OCR 정확도 영향 미미
_OCR_MAX_SIDE = 1920         # FHD 유지 (작은 글자 위해 너무 줄이면 안 됨)


def set_encode_quality(jpeg_quality: int = 90, max_side: int = 1920):
    """
    속도 조정용 전역 설정.
    - jpeg_quality: 낮을수록 업로드 용량↓ (권장 80~95)
    - max_side   : 이미지 최대 변. 1920 미만으로 낮추면 작은 글자 판독 저하 주의.
    """
    global _JPEG_QUALITY, _OCR_MAX_SIDE
    _JPEG_QUALITY = max(60, min(100, int(jpeg_quality)))
    _OCR_MAX_SIDE = max(720, int(max_side))


def _preprocess_for_ocr(image):
    """OCR 전 이미지 전처리: 리사이즈 + CLAHE + 언샤프."""
    h, w = image.shape[:2]
    if max(w, h) > _OCR_MAX_SIDE:
        scale = _OCR_MAX_SIDE / max(w, h)
        image = cv2.resize(image, (int(w * scale), int(h * scale)),
                           interpolation=cv2.INTER_LANCZOS4)

    lab = cv2.cvtColor(image, cv2.COLOR_BGR2LAB)
    l, a, b = cv2.split(lab)
    clahe = cv2.createCLAHE(clipLimit=2.5, tileGridSize=(8, 8))
    l = clahe.apply(l)
    image = cv2.cvtColor(cv2.merge([l, a, b]), cv2.COLOR_LAB2BGR)

    gaussian = cv2.GaussianBlur(image, (0, 0), 2.0)
    image = cv2.addWeighted(image, 1.5, gaussian, -0.5, 0)

    return image


def encode_image(image):
    """전처리된 이미지를 JPEG Base64로 인코딩합니다."""
    image = _preprocess_for_ocr(image)
    _, buffer = cv2.imencode('.jpg', image, [int(cv2.IMWRITE_JPEG_QUALITY), _JPEG_QUALITY])
    return base64.b64encode(buffer).decode('utf-8')


def _phash(image_bgr, hash_size: int = 16) -> int:
    """
    간단한 perceptual hash (DCT 기반).
    자막 등 세밀한 변화까지 감지하도록 hash_size를 16으로 상향.
    """
    gray = cv2.cvtColor(image_bgr, cv2.COLOR_BGR2GRAY)
    resized = cv2.resize(gray, (hash_size * 4, hash_size * 4),
                         interpolation=cv2.INTER_AREA).astype(np.float32)
    dct = cv2.dct(resized)
    dct_low = dct[:hash_size, :hash_size]
    # DC 성분(좌상단) 제외 평균
    med = np.median(dct_low.flatten()[1:])
    bits = (dct_low > med).flatten()
    h = 0
    for b in bits:
        h = (h << 1) | int(b)
    return h


def _hamming(a: int, b: int) -> int:
    return bin(a ^ b).count("1")


def _focus_region(gray: np.ndarray) -> np.ndarray:
    """
    프레임의 좌측상단(챕터/소주제) + 우측상단(STEP 로고) + 강사 영역을 제외한
    '중앙 콘텐츠 노란 박스 영역'만 잘라 반환합니다. 안정성 판정에 사용 — 3D 배경의
    아주 작은 흔들림이나 강사 미세 동작이 전체 차이값을 키우는 것을 방지.

    입력 gray는 (180, 320) 또는 비슷한 작은 크기의 그레이스케일 numpy 배열.
    """
    h, w = gray.shape[:2]
    # 좌측상단 25%·우측상단 18%·하단 강사 영역(아래 22%) 마스킹 후 중앙 영역
    y1 = int(h * 0.18)    # 상단 헤더 잘라냄
    y2 = int(h * 0.82)    # 하단 강사 잘라냄
    x1 = int(w * 0.05)
    x2 = int(w * 0.95)
    if y2 - y1 < 20 or x2 - x1 < 20:
        return gray
    return gray[y1:y2, x1:x2]


def _extract_with_stability(
    video_path: str,
    check_interval: float = 0.35,
    motion_threshold: float = 1.5,
    min_stable_seconds: float = 1.8,
    capture_min_gap: float = 1.5,
    phash_threshold: int = 10,
) -> List[dict]:
    """
    안정화 감지 기반 프레임 추출 — **모션 종료 프레임 우선** 버전.

    기존 방식은 화면이 안정 상태에 처음 도달한 순간을 캡처했습니다.
    하지만 교육 영상은 텍스트/도형이 단계적으로 등장하면서 중간중간 잠깐 멈추는 경우가 많아
    스토리보드의 완성 화면이 아니라 애니메이션 중간 화면이 잡힐 수 있습니다.

    이 버전은 안정 상태가 시작되면 바로 저장하지 않고 후보 프레임만 계속 갱신하다가,
    다음 모션이 시작되거나 영상이 끝날 때 **직전 안정 구간의 마지막 프레임**을 저장합니다.
    즉, 한 장면의 모션이 모두 끝난 상태를 스토리보드와 비교합니다.
    """
    cap = cv2.VideoCapture(video_path)
    fps = cap.get(cv2.CAP_PROP_FPS) or 30.0
    interval_frames = max(1, int(fps * check_interval))
    min_stable_count = max(2, int(round(min_stable_seconds / max(check_interval, 0.05))))
    min_gap_frames = max(1, int(fps * capture_min_gap))

    unique_frames: List[dict] = []
    prev_hashes: List[int] = []
    prev_focus_gray: Optional[np.ndarray] = None
    stable_streak: int = 0
    stable_candidate_frame: Optional[np.ndarray] = None
    stable_candidate_idx: int = -1
    last_capture_frame_idx: int = -10**9
    last_full_frame: Optional[np.ndarray] = None
    last_frame_idx: int = -1

    def _try_emit_candidate():
        """현재 안정 후보의 마지막 프레임을 중복 제거 후 결과에 추가."""
        nonlocal stable_candidate_frame, stable_candidate_idx, last_capture_frame_idx
        if stable_candidate_frame is None or stable_candidate_idx < 0:
            return
        if stable_candidate_idx - last_capture_frame_idx < min_gap_frames:
            return
        h_now = _phash(stable_candidate_frame)
        is_dup = any(_hamming(h_now, ph) <= phash_threshold for ph in prev_hashes[-10:])
        if is_dup:
            return
        current_sec = stable_candidate_idx / fps
        unique_frames.append({
            "time": current_sec,
            "time_str": format_timestamp(current_sec),
            "base64": encode_image(stable_candidate_frame),
        })
        prev_hashes.append(h_now)
        last_capture_frame_idx = stable_candidate_idx

    frame_count = 0
    while True:
        if not cap.grab():
            break
        if frame_count % interval_frames != 0:
            frame_count += 1
            continue

        ret, frame = cap.retrieve()
        if not ret:
            frame_count += 1
            continue

        small = cv2.resize(frame, (320, 180))
        gray = cv2.cvtColor(small, cv2.COLOR_BGR2GRAY)
        focus = _focus_region(gray)

        if prev_focus_gray is not None:
            diff = float(np.mean(cv2.absdiff(focus, prev_focus_gray)))
            is_static = (diff < motion_threshold)
        else:
            is_static = False

        if is_static:
            stable_streak += 1
            if (stable_streak + 1) >= min_stable_count:
                stable_candidate_frame = frame.copy()
                stable_candidate_idx = frame_count
        else:
            if stable_candidate_frame is not None:
                _try_emit_candidate()
            stable_streak = 0
            stable_candidate_frame = None
            stable_candidate_idx = -1

        prev_focus_gray = focus
        last_full_frame = frame
        last_frame_idx = frame_count
        frame_count += 1

    if stable_candidate_frame is not None:
        _try_emit_candidate()
    elif last_full_frame is not None and stable_streak > 0:
        stable_candidate_frame = last_full_frame
        stable_candidate_idx = last_frame_idx
        _try_emit_candidate()

    cap.release()
    return unique_frames

def extract_and_filter_frames(
    video_path,
    sample_rate: float = 1.0,
    diff_threshold: float = 15.0,
    phash_threshold: int = 10,
    scene_change_threshold: float = 60.0,
    # ★ v4: 안정화(애니메이션 종료) 감지 모드 ─ 기본 ON
    stability_mode: bool = True,
    stability_check_interval: float = 0.35,  # 내부 샘플링 간격 (초) - 안정성 추적용
    stability_motion_threshold: float = 1.5,  # 인접 샘플 간 평균픽셀차 ≤ 이면 "정지"
    stability_min_seconds: float = 1.8,       # 이 시간 이상 정지 지속 시 안정으로 판정
    capture_min_gap: float = 1.5,             # 같은 장면 중복 캡처 방지: 마지막 캡처 후 최소 간격
):
    """
    영상에서 sample_rate 간격으로 프레임을 추출하고
    (1) 저해상도 평균차이로 빠른 필터링
    (2) pHash 해밍거리로 정밀 중복 제거

    ★ 속도 업그레이드:
    - `cap.grab()` + 조건부 `retrieve()` 패턴: 샘플링 간격에 해당하지 않는
      프레임은 디코딩 자체를 건너뛴다 → 대부분의 영상에서 2~3배 빠름.
    - 전체 프레임 수를 미리 계산해 진행 예측 가능.

    ★ 정확도 v2:
    - scene_change_threshold: 샘플링 간격 사이에서도 급격한 변화(장면 전환)가
      감지되면 해당 프레임을 추가로 캡처. 스토리보드 슬라이드 전환을 놓치지 않음.

    ★ v4 — 안정화(애니메이션 완료 후) 캡처 모드:
    - 영상에 페이드인·도형 그리기·텍스트 등장 같은 모션이 있으면 기존 알고리즘은
      애니메이션 중간 프레임을 캡처해 스토리보드(완성 상태)와 비교 시 위양성 오류가 발생.
    - stability_mode=True 인 경우, 짧은 간격(stability_check_interval)으로
      프레임을 살피고 인접 프레임 간 픽셀차가 stability_motion_threshold 미만으로
      stability_min_seconds 이상 지속되면 "안정 상태"로 간주하여 그 마지막 프레임만 캡처.
    - 이렇게 하면 애니메이션이 끝나고 완성된 화면이 잠시 유지되는 시점에서만
      캡처가 일어나 스토리보드 매칭 정확도가 크게 올라감.
    - 동일 장면을 여러 번 잡지 않도록 capture_min_gap 시간 안의 중복은 pHash로 차단.

    phash_threshold: 해밍거리 이하는 중복으로 간주 (기본 10 / 256비트 기준)
    """
    if stability_mode:
        return _extract_with_stability(
            video_path,
            check_interval=stability_check_interval,
            motion_threshold=stability_motion_threshold,
            min_stable_seconds=stability_min_seconds,
            capture_min_gap=capture_min_gap,
            phash_threshold=phash_threshold,
        )
    # ── 이하 기존 sample_rate 기반 알고리즘 (호환성 유지) ─────────────
    cap = cv2.VideoCapture(video_path)
    fps = cap.get(cv2.CAP_PROP_FPS) or 30.0
    frame_interval = max(1, int(fps * sample_rate))
    # 장면 전환 감지용 서브샘플 간격 (0.2초마다 가벼운 체크)
    scene_check_interval = max(1, int(fps * 0.2))

    unique_frames = []
    prev_gray = None
    prev_hashes: List[int] = []
    prev_scene_gray = None  # 장면 전환 감지용
    frame_count = 0

    while True:
        if not cap.grab():
            break

        is_sample_point = (frame_count % frame_interval == 0)
        is_scene_check = (frame_count % scene_check_interval == 0) and not is_sample_point

        should_decode = is_sample_point or is_scene_check

        if should_decode:
            ret, frame = cap.retrieve()
            if not ret:
                frame_count += 1
                continue

            small = cv2.resize(frame, (320, 180))
            gray = cv2.cvtColor(small, cv2.COLOR_BGR2GRAY)

            if is_scene_check and not is_sample_point:
                # 장면 전환 감지만 수행 (가벼운 체크)
                if prev_scene_gray is not None:
                    scene_diff = np.mean(cv2.absdiff(gray, prev_scene_gray))
                    if scene_diff >= scene_change_threshold:
                        # 급격한 변화 감지 → 이 프레임을 유니크로 추가
                        h = _phash(frame)
                        if not any(_hamming(h, ph) <= phash_threshold for ph in prev_hashes[-5:]):
                            current_sec = frame_count / fps
                            unique_frames.append({
                                "time": current_sec,
                                "time_str": format_timestamp(current_sec),
                                "base64": encode_image(frame),
                            })
                            prev_gray = gray
                            prev_hashes.append(h)
                prev_scene_gray = gray
                frame_count += 1
                continue

            # 정규 샘플링 포인트 처리
            is_unique = True
            if prev_gray is not None:
                diff = cv2.absdiff(gray, prev_gray)
                if np.mean(diff) < diff_threshold:
                    h = _phash(frame)
                    if any(_hamming(h, ph) <= phash_threshold for ph in prev_hashes[-5:]):
                        is_unique = False
                    else:
                        prev_hashes.append(h)
                else:
                    prev_hashes.append(_phash(frame))
            else:
                prev_hashes.append(_phash(frame))

            if is_unique:
                current_sec = frame_count / fps
                unique_frames.append({
                    "time": current_sec,
                    "time_str": format_timestamp(current_sec),
                    "base64": encode_image(frame),
                })
                prev_gray = gray

            prev_scene_gray = gray

        frame_count += 1

    cap.release()
    return unique_frames


# ─────────────────────────────────────────────
# 스토리보드 PPT 이미지 추출 + 메타데이터 파싱
# ─────────────────────────────────────────────

def _natural_sort_key(s):
    return [int(t) if t.isdigit() else t.lower() for t in re.split(r'(\d+)', s)]


# ─────────────────────────────────────────────
# PPT 슬라이드 영역 분류 상수
# ─────────────────────────────────────────────
# 슬라이드를 5개 영역으로 분할 (정규화 좌표 0.0~1.0):
#   ┌──────── HEADER (top < 0.10) ────────┐
#   │INDEX │     CENTER     │ R-DESC      │
#   │(0.13)│  (영상 비교)    │ (0.85~)     │
#   │      │                 │             │
#   ├──────┴─────────────────┴─────────────┤
#   │     NARRATION (top > 0.78)           │
#   └──────────────────────────────────────┘
_REGION_HEADER_TOP = 0.10
_REGION_NARRATION_TOP = 0.78
_REGION_INDEX_LEFT = 0.13
_REGION_RDESC_LEFT = 0.85


def _classify_shape_region(left_n: float, top_n: float) -> str:
    """
    정규화된 (left, top) 좌표로부터 슬라이드 영역을 분류합니다.
    Returns: "header" | "index" | "rdesc" | "narration" | "center"
    """
    if top_n >= _REGION_NARRATION_TOP:
        return "narration"
    if top_n < _REGION_HEADER_TOP:
        return "header"
    if left_n < _REGION_INDEX_LEFT:
        return "index"
    if left_n >= _REGION_RDESC_LEFT:
        return "rdesc"
    return "center"


# ─────────────────────────────────────────────
# ★ v11/v12: 메타 슬라이드 자동 판별 (키워드 기반)
# ─────────────────────────────────────────────
# v12 변경: 키워드를 매우 엄격하게 정의 — 콘텐츠 슬라이드에 우연히 등장할 수 있는
# 짧은 일반 단어("용어집", "표시사항" 등)는 제거하고, 명확한 메타 헤더만 유지.
# 또한 메타 판별 이유를 info["meta_reason"] 에 기록해서 사용자가 확인 가능.

# 정확한 메타 슬라이드 헤더 키워드 (이런 문구가 슬라이드 상단에 오면 100% 메타)
_META_HEADER_KEYWORDS = (
    "스토리보드 설계 가이드",
    "콘텐츠 용어집",
    "스토리보드 표시사항",
    "스토리보드 가이드",
    "[과정정보]",
    "[NCS 정보]",
    "[NCS정보]",
)

# 메타 슬라이드에서만 함께 등장하는 강력한 보조 키워드들
# (단독 매칭이 아닌, 여러 개 동시에 보여야 메타로 판정)
_META_AUX_KEYWORDS = (
    "내용전문가",
    "회차 키워드",
    "수정 사유",
    "대치 필요 용어",
    "원고집필자",
    "내용 검토 자문가",
    "교수자 정보",
    "문의/요청",
    "스토리보드 표시",
)


def _detect_meta_slide(info: dict) -> bool:
    """
    슬라이드가 메타 슬라이드인지 키워드 기반으로 판별.

    ★ v12: 매우 엄격한 매칭으로 변경 (false positive 방지)
      1) _META_HEADER_KEYWORDS 중 하나라도 정확히 포함 → 메타
      2) _META_AUX_KEYWORDS 중 2개 이상 동시 등장 → 메타
      3) 슬라이드가 완전 비어있음 (텍스트 0개) → 메타

    이전의 "텍스트 적음 → 메타" 규칙은 큰 이미지 콘텐츠 슬라이드를 잘못
    분류할 수 있어 제거되었습니다.

    Returns: True면 메타. info["meta_reason"]에 판별 이유 기록.
    """
    info["meta_reason"] = ""
    all_texts = info.get("all_texts", []) or []

    # 0) 완전 빈 슬라이드만 메타로 처리 (1개 이상이면 콘텐츠로 간주)
    if not all_texts:
        info["meta_reason"] = "슬라이드가 비어 있음"
        return True

    # 정규화된 텍스트
    text_blob = " ".join(t.replace(" ", "") for t in all_texts).lower()

    # 1) 정확한 메타 헤더 키워드 매칭 (단독으로 강한 신호)
    for kw in _META_HEADER_KEYWORDS:
        kw_norm = kw.replace(" ", "").lower()
        if kw_norm in text_blob:
            info["meta_reason"] = f"메타 헤더 '{kw}' 감지"
            return True

    # 2) 보조 키워드가 2개 이상 동시에 등장하는 경우만 메타
    aux_hits = []
    for kw in _META_AUX_KEYWORDS:
        kw_norm = kw.replace(" ", "").lower()
        if kw_norm in text_blob:
            aux_hits.append(kw)
    if len(aux_hits) >= 2:
        info["meta_reason"] = f"보조 키워드 다수 감지: {', '.join(aux_hits[:3])}"
        return True

    return False


def extract_ppt_metadata(ppt_path: str) -> List[dict]:
    """
    python-pptx로 스토리보드 PPT의 메타데이터를 추출합니다.

    각 슬라이드에서 다음을 파싱:
    - screen_no: 화면번호 (예: "03_01", "03_02")
    - narration: 내레이션 텍스트 (하단 영역의 "교수:" 텍스트)
    - title: 화면 제목
    - is_meta: 메타 슬라이드 여부 (용어집, 설계 가이드 등)
    - is_title_screen: 섹션 타이틀 화면 여부 (XX_01 패턴)
    - all_texts: 슬라이드 내 모든 텍스트 (호환용)
    - region_texts: 영역별 분류된 텍스트 (위치 기반)

    ★ v9 추가:
    - 그룹(GROUP) 및 중첩 shape를 재귀적으로 탐색
    - shape의 left/top이 None인 경우 부모 좌표로 폴백
    - 통계 정보 print로 추출 품질 확인 가능
    """
    try:
        from pptx import Presentation as PptxPresentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print("python-pptx가 설치되어 있지 않습니다. 메타데이터 없이 진행합니다.")
        return []

    try:
        prs = PptxPresentation(ppt_path)
    except Exception as e:
        print(f"PPT 메타데이터 파싱 오류: {e}")
        return []

    slide_w = prs.slide_width or 9144000
    slide_h = prs.slide_height or 6858000

    _SCREEN_NO_RE = re.compile(r'^\d+_\d+(?:_\d+)?$')   # ★ v11: 자릿수 유연 (01_01, 1_2 등 모두 지원)
    _SIMPLE_NO_RE = re.compile(r'^\d{1,3}$')              # ★ v11: 단순 페이지 번호 ("01", "1", "001")
    _NARRATION_RE = re.compile(r'교수\s*:')

    def _iter_text_shapes(shape, parent_left=0, parent_top=0):
        """
        ★ v9: 재귀적으로 모든 텍스트 shape를 평탄화해서 yield합니다.
        그룹 도형, 중첩 그룹까지 탐색. 좌표가 None인 경우 부모 좌표로 폴백.

        Yields: (shape, abs_left, abs_top)
        """
        try:
            sl = shape.left
            st_ = shape.top
        except Exception:
            sl, st_ = None, None

        # 좌표 None 폴백 → 부모 좌표 사용
        if sl is None:
            sl = parent_left
        if st_ is None:
            st_ = parent_top

        # 그룹이면 자식들 재귀 탐색
        try:
            shape_type = shape.shape_type
        except Exception:
            shape_type = None

        if shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                children = list(shape.shapes)
            except Exception:
                children = []
            for sub in children:
                yield from _iter_text_shapes(sub, parent_left=sl, parent_top=st_)
            return

        # 텍스트 또는 테이블이 있는 일반 shape
        yield (shape, sl, st_)

    metadata = []
    extraction_stats = {"total": 0, "with_center": 0, "empty_shapes_in_slides": []}

    for slide_idx, slide in enumerate(prs.slides):
        info = {
            "slide_num": slide_idx + 1,
            "screen_no": "",
            "narration": "",
            "title": "",
            "is_meta": False,
            "meta_reason": "",   # ★ v12: 메타 판별 이유 (디버깅용)
            "is_title_screen": False,
            "all_texts": [],
            "region_texts": {
                "center": [], "narration": [],
                "header": [], "index": [], "rdesc": [],
            },
        }

        texts = []

        def _add_text_with_region(t: str, left, top):
            if not t:
                return
            t = t.strip()
            if not t:
                return
            texts.append(t)
            try:
                left_n = (left or 0) / slide_w
                top_n = (top or 0) / slide_h
            except Exception:
                left_n, top_n = 0.5, 0.5
            region = _classify_shape_region(left_n, top_n)
            info["region_texts"][region].append(t)

            if _SCREEN_NO_RE.match(t) and not info["screen_no"]:
                info["screen_no"] = t
            # ★ v11: NN_NN 형태가 못 찾을 때, 우측상단 영역의 짧은 숫자도 인정
            #   (모의해킹 PPT처럼 "01"만 있는 양식 대응)
            elif (not info["screen_no"]
                  and _SIMPLE_NO_RE.match(t)
                  and 1 <= len(t) <= 3
                  and region == "header"
                  and left_n > 0.85):     # 우측상단 영역
                info["screen_no"] = t

            if region == "narration" and _NARRATION_RE.search(t) and not info["narration"]:
                clean = re.sub(r'#\d+\s*', '', t)
                clean = re.sub(r'교수\s*:\s*', '', clean).strip()
                info["narration"] = clean

        # ★ v9: 재귀적으로 모든 shape 평탄화해서 순회
        for shape, abs_l, abs_t in _iter_text_shapes(slide.shapes):
            try:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        _add_text_with_region(t, abs_l, abs_t)
                elif shape.has_table:
                    table = shape.table
                    rows = list(table.rows)
                    cols = list(table.columns)
                    n_rows = len(rows) or 1
                    n_cols = len(cols) or 1
                    shape_w = shape.width or 0
                    shape_h = shape.height or 0
                    base_l = abs_l or 0
                    base_t = abs_t or 0
                    for ri, row in enumerate(rows):
                        for ci, cell in enumerate(row.cells):
                            ct = cell.text.strip()
                            if not ct:
                                continue
                            cell_l = base_l + (shape_w * ci / n_cols)
                            cell_t = base_t + (shape_h * ri / n_rows)
                            _add_text_with_region(ct, cell_l, cell_t)
            except Exception:
                continue

        info["all_texts"] = texts

        # ★ v11: 메타 슬라이드 판별 — 키워드 기반 (screen_no 의존성 제거)
        # 다양한 PPT 양식에서 모두 작동하도록, 슬라이드 내용으로 판별합니다.
        info["is_meta"] = _detect_meta_slide(info)

        # is_title_screen: screen_no 있을 때만 판별
        if info["screen_no"]:
            parts = info["screen_no"].split("_")
            if len(parts) >= 2 and parts[1] in ("01", "1"):
                info["is_title_screen"] = True

        for src_region in ("index", "header", "center"):
            for t in info["region_texts"].get(src_region, []):
                if (t != info["screen_no"]
                    and not _NARRATION_RE.search(t)
                    and 2 < len(t) < 40
                    and t not in ("자막영역 침범금지", "BGM", "내용 제시", "INDEX",
                                  "과정명", "회차명", "화면번호", "화면설명",
                                  "이미지 번호", "내용을 입력하십시오")
                    and not t.startswith("#")
                    and not t.startswith("그림")):
                    info["title"] = t
                    break
            if info["title"]:
                break

        # 통계 수집
        extraction_stats["total"] += 1
        if info["region_texts"]["center"]:
            extraction_stats["with_center"] += 1
        else:
            # ALLSLIDES v2: 메타로 감지된 슬라이드도 제외하지 않고 품질 진단에 포함
            extraction_stats["empty_shapes_in_slides"].append(info["slide_num"])

        metadata.append(info)

    # ★ v9: 추출 품질 진단 출력
    n_total = extraction_stats["total"]
    n_with_center = extraction_stats["with_center"]
    n_empty = len(extraction_stats["empty_shapes_in_slides"])
    print(f"[PPT 추출] 전체 {n_total}장 중 중앙 텍스트 추출 성공: {n_with_center}장")
    if n_empty > 0:
        empty_list = extraction_stats["empty_shapes_in_slides"][:10]
        print(f"  ⚠️ 중앙 텍스트 추출 실패 (콘텐츠 슬라이드): {n_empty}장 — "
              f"슬라이드 번호: {empty_list}{' ...' if n_empty > 10 else ''}")
        print(f"  → 이 슬라이드들은 중앙 콘텐츠가 PPT에 이미지로 박혀있을 가능성. "
              f"Vision OCR 폴백을 활성화하면 자동 추출 가능.")

    return metadata


def extract_ppt_metadata_LEGACY_(ppt_path: str) -> List[dict]:
    """(보존용) 영역 분류 없는 구버전 — 더 이상 사용하지 않음."""
    return []   # placeholder


def _crop_region_b64(
    img_b64: str,
    left_ratio: float,
    top_ratio: float,
    right_ratio: float,
    bottom_ratio: float,
    *,
    label: str = "영역",
    max_size: Tuple[int, int] = (1280, 720),
) -> str:
    """Base64 이미지를 비율 좌표로 크롭합니다. 실패 시 원본을 반환합니다."""
    try:
        raw = base64.b64decode(img_b64)
        pil = PILImage.open(io.BytesIO(raw)).convert("RGB")
        W, H = pil.size
        l = max(0, min(W - 1, int(W * left_ratio)))
        t = max(0, min(H - 1, int(H * top_ratio)))
        r = max(l + 1, min(W, int(W * right_ratio)))
        b = max(t + 1, min(H, int(H * bottom_ratio)))
        cropped = pil.crop((l, t, r, b))
        cropped.thumbnail(max_size, PILImage.LANCZOS)
        buf = io.BytesIO()
        cropped.save(buf, format="JPEG", quality=88)
        buf.seek(0)
        return base64.b64encode(buf.read()).decode("utf-8")
    except Exception as e:
        print(f"{label} 크롭 실패: {e}")
        return img_b64


def crop_content_region_b64(sb_image_b64: str) -> str:
    """
    스토리보드에서 영상 화면과 비교해야 하는 **중앙 콘텐츠 노란 박스**만 잘라냅니다.
    좌측 INDEX, 상단 표 헤더, 우측 화면설명/이미지번호, 하단 내레이션 영역을 제외합니다.
    """
    return _crop_region_b64(
        sb_image_b64,
        _REGION_INDEX_LEFT, _REGION_HEADER_TOP,
        _REGION_RDESC_LEFT, _REGION_NARRATION_TOP,
        label="중앙 콘텐츠 노란 박스 영역",
        max_size=(1280, 720),
    )


def crop_narration_region_b64(sb_image_b64: str) -> str:
    """
    스토리보드에서 음성과 비교해야 하는 **하단 내레이션 노란 박스**만 잘라냅니다.
    실제 음성 비교는 PPT에서 추출한 내레이션 텍스트를 기준으로 수행하고,
    이 이미지는 OCR 보강/디버깅에 사용할 수 있습니다.
    """
    return _crop_region_b64(
        sb_image_b64,
        0.0, _REGION_NARRATION_TOP,
        _REGION_RDESC_LEFT, 1.0,
        label="하단 내레이션 영역",
        max_size=(1280, 260),
    )


# 호환용 alias (구 명칭): 기존 코드가 crop_center_region_b64/crop_slide_centers를 호출해도
# 실제로는 노란색 중앙 콘텐츠 박스를 사용합니다.
crop_center_region_b64 = crop_content_region_b64
_crop_center_region_b64 = crop_content_region_b64


def crop_slide_content_regions(sb_images: List[str]) -> List[str]:
    """모든 스토리보드 슬라이드의 중앙 콘텐츠 노란 박스만 크롭합니다."""
    return [crop_content_region_b64(img) for img in sb_images]


def crop_slide_narration_regions(sb_images: List[str]) -> List[str]:
    """모든 스토리보드 슬라이드의 하단 내레이션 노란 박스만 크롭합니다."""
    return [crop_narration_region_b64(img) for img in sb_images]


def crop_slide_centers(sb_images: List[str]) -> List[str]:
    """기존 이름 호환용. 실제 반환값은 중앙 콘텐츠 노란 박스 크롭 이미지입니다."""
    return crop_slide_content_regions(sb_images)


def crop_video_content_region_b64(frame_b64: str) -> str:
    """
    영상 프레임에서 스토리보드의 중앙 콘텐츠 노란 박스와 대응되는 영역만 잘라냅니다.

    - 상단 좌측 2줄(차시명/소주제)과 우측 STEP 로고 영역은 별도 검증 대상이므로
      이미지 매칭에서는 제외합니다.
    - 하단 자막/내레이션 영역도 음성 비교 대상이므로 제외합니다.
    - 템플릿마다 여백이 조금 달라도 대응되도록 살짝 넓게 자릅니다.
    """
    return _crop_region_b64(
        frame_b64,
        0.02, 0.10,
        0.98, 0.88,
        label="영상 중앙 콘텐츠 영역",
        max_size=(1280, 720),
    )


def crop_video_header_left_region_b64(frame_b64: str) -> str:
    """영상 좌측상단 2줄(차시명/소주제) 검증용 크롭."""
    return _crop_region_b64(frame_b64, 0.00, 0.00, 0.55, 0.14, label="영상 좌측상단 2줄", max_size=(720, 140))


def crop_video_step_logo_region_b64(frame_b64: str) -> str:
    """영상 우측상단 STEP 로고 검증용 크롭."""
    return _crop_region_b64(frame_b64, 0.70, 0.00, 1.00, 0.16, label="영상 우측상단 STEP 로고", max_size=(480, 160))


_SLIDE_OCR_PROMPT = """이 이미지는 교육용 스토리보드 PPT 슬라이드의 **중앙 콘텐츠 노란 박스 영역**입니다.
이 영역에 보이는 모든 한국어/영어 텍스트(제목, 라벨, 본문, 박스 안 글자, 화살표 옆 글자 등)를
빠짐없이 추출해 JSON 배열로 반환하세요.

요구사항:
- 한 줄에 하나씩, 자연스러운 단위로 분리 (한 도형 안의 여러 줄은 합쳐도 됨)
- 의미 없는 단순 기호(•, ■ 같은 글머리표)는 제외
- 페이지 번호, "그림 N" 같은 보일러플레이트는 제외
- 결과 형식: {"texts": ["텍스트1", "텍스트2", ...]}

만약 텍스트가 전혀 없으면 {"texts": []} 반환.
"""

_SLIDE_OCR_SCHEMA = {
    "name": "slide_ocr_result",
    "strict": True,
    "schema": {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "texts": {"type": "array", "items": {"type": "string"}},
        },
        "required": ["texts"],
    },
}


def vision_ocr_slide_center(
    sb_image_b64: str,
    api_key: str,
    model: str = "gpt-5.4-mini",
) -> List[str]:
    """
    스토리보드 슬라이드 이미지의 중앙 영역을 GPT Vision으로 OCR하여 텍스트 리스트 반환.
    PPT shape에서 텍스트 추출이 실패한 슬라이드(이미지로 박힌 텍스트)에 사용.
    """
    try:
        cropped_b64 = _crop_center_region_b64(sb_image_b64)
        client = OpenAI(api_key=api_key)
        response = call_with_retry(
            lambda: client.chat.completions.create(
                model=model,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": _SLIDE_OCR_PROMPT},
                        {"type": "image_url", "image_url": {
                            "url": f"data:image/jpeg;base64,{cropped_b64}",
                            "detail": "high",
                        }},
                    ],
                }],
                temperature=0.0,
                response_format={"type": "json_schema", "json_schema": _SLIDE_OCR_SCHEMA},
            )
        )
        parsed = json.loads(response.choices[0].message.content)
        texts = parsed.get("texts", []) or []
        # 정제
        out = []
        seen = set()
        for t in texts:
            tt = (t or "").strip()
            if not tt or len(tt) < 2 or len(tt) > 200:
                continue
            if tt in seen:
                continue
            seen.add(tt)
            out.append(tt)
        return out
    except Exception as e:
        print(f"Vision OCR 실패: {e}")
        return []


def enrich_slide_metadata_with_ocr(
    metadata: List[dict],
    sb_images: List[str],
    api_key: str,
    *,
    model: str = "gpt-5.4-mini",
    only_empty_center: bool = True,
    max_workers: int = 4,
    progress_cb=None,
) -> int:
    """
    metadata의 region_texts['center']가 비어있는 슬라이드에 대해 Vision OCR을
    수행하고 결과를 region_texts['center']와 all_texts에 추가합니다.

    Parameters
    ----------
    only_empty_center : True면 중앙이 빈 슬라이드만 OCR (비용 절약)
                       False면 모든 슬라이드에 대해 OCR 수행 (보강)
    progress_cb : callable(done, total) 진행률 콜백 (선택)

    Returns
    -------
    int : OCR 적용된 슬라이드 수
    """
    if not metadata or not sb_images:
        return 0

    targets = []
    for i, info in enumerate(metadata):
        if i >= len(sb_images):
            break
        # ALLSLIDES v2: 메타 감지 여부와 무관하게 모든 슬라이드 OCR 보강 대상에 포함
        center = info.get("region_texts", {}).get("center", []) or []
        if only_empty_center and len(center) >= 2:
            continue   # 이미 충분한 텍스트 있음
        targets.append(i)

    if not targets:
        return 0

    print(f"[Vision OCR 폴백] 중앙 텍스트가 빈약한 {len(targets)}장 슬라이드를 OCR로 보강합니다...")
    done = 0
    total = len(targets)

    def _process_one(idx):
        ocr_texts = vision_ocr_slide_center(sb_images[idx], api_key, model=model)
        return idx, ocr_texts

    enriched_count = 0
    with ThreadPoolExecutor(max_workers=max_workers) as ex:
        futures = {ex.submit(_process_one, i): i for i in targets}
        for fut in as_completed(futures):
            try:
                idx, ocr_texts = fut.result()
                if ocr_texts:
                    info = metadata[idx]
                    region = info.setdefault("region_texts", {})
                    center_list = region.setdefault("center", [])
                    seen = set(center_list)
                    for t in ocr_texts:
                        if t not in seen:
                            seen.add(t)
                            center_list.append(t)
                            info.setdefault("all_texts", []).append(t)
                    enriched_count += 1
            except Exception as e:
                print(f"OCR 실패 (슬라이드 {idx + 1}): {e}")
            done += 1
            if progress_cb:
                try:
                    progress_cb(done, total)
                except Exception:
                    pass

    print(f"[Vision OCR 폴백] 완료: {enriched_count}/{total}장 보강.")
    return enriched_count


def extract_ppt_slides(ppt_path, output_dir):
    """PPT 파일을 JPG로 내보낸 후 base64로 인코딩한 이미지 리스트 반환"""
    sb_images = []
    powerpoint = None
    try:
        pythoncom.CoInitialize()
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(os.path.abspath(ppt_path), WithWindow=False)
        presentation.Export(os.path.abspath(output_dir), "JPG")
        presentation.Close()

        files = sorted(
            [f for f in os.listdir(output_dir) if f.lower().endswith('.jpg')],
            key=_natural_sort_key
        )
        for f in files:
            with open(os.path.join(output_dir, f), "rb") as im:
                sb_images.append(base64.b64encode(im.read()).decode("utf-8"))
    except Exception as e:
        print(f"PPT 추출 오류: {e}")
    finally:
        try:
            if powerpoint and powerpoint.Presentations.Count == 0:
                powerpoint.Quit()
        except:
            pass
        pythoncom.CoUninitialize()
    return sb_images


def get_storyboard_structure(metadata: List[dict]) -> dict:
    """
    스토리보드 메타데이터로부터 영상별 슬라이드 매핑 구조를 생성합니다.

    ★ v2 변경: 각 slide에 chapter/subtopic을 역으로 전파합니다.
      - chapter  = 해당 슬라이드가 속한 섹션(XX)의 타이틀 슬라이드(XX_01)의 title
      - subtopic = 슬라이드 자체의 title (타이틀 슬라이드면 chapter와 동일)
    이 정보는 영상의 좌측상단 2줄(챕터/소주제) 검증에 사용됩니다.

    Returns
    -------
    {
        "meta_slides": [1,2,3,4,5,6],
        "content_slides": [7,8,9,...],
        "sections": {
            "03": {
                "title": "TSV 공정기술",
                "title_slide": 11,
                "content_slides": [12,13,...,28],
                "screen_nos": ["03_01", "03_02"],
            }, ...
        },
        "narrations": [...],
        "screen_no_to_slides": {...},
    }
    """
    result = {
        "meta_slides": [],
        "content_slides": [],
        "sections": {},
        "narrations": [],
        "screen_no_to_slides": {},
    }

    for info in metadata:
        sn = info["slide_num"]
        if info["is_meta"]:
            result["meta_slides"].append(sn)
        else:
            result["content_slides"].append(sn)

        # 화면번호별 슬라이드 그룹핑
        scr = info["screen_no"]
        if scr:
            result["screen_no_to_slides"].setdefault(scr, []).append(sn)

            # 섹션별 정리
            sec = scr.split("_")[0]
            if sec not in result["sections"]:
                result["sections"][sec] = {
                    "title": "",
                    "title_slide": None,
                    "content_slides": [],
                    "screen_nos": set(),
                }
            result["sections"][sec]["screen_nos"].add(scr)

            if info["is_title_screen"]:
                result["sections"][sec]["title_slide"] = sn
                result["sections"][sec]["title"] = info["title"]
            else:
                result["sections"][sec]["content_slides"].append(sn)

        # 내레이션 수집
        if info["narration"]:
            result["narrations"].append(info["narration"])

    # sections의 screen_nos를 리스트로 변환 (JSON 직렬화용)
    for sec in result["sections"].values():
        sec["screen_nos"] = sorted(sec["screen_nos"])

    # ★ v2: chapter/subtopic을 각 슬라이드 metadata에 역으로 주입
    #   (metadata는 외부 참조(mutable)를 그대로 사용하므로 여기서 in-place 수정)
    for info in metadata:
        scr = info.get("screen_no") or ""
        sec = scr.split("_")[0] if scr else ""
        section_info = result["sections"].get(sec) if sec else None
        # chapter: 해당 섹션의 타이틀
        chapter = section_info["title"] if section_info else ""
        # subtopic: 해당 슬라이드의 title — 없으면 chapter로 대체
        subtopic = info.get("title") or chapter
        info["chapter"] = chapter
        info["subtopic"] = subtopic

    return result


def _extract_content_texts(info: dict, max_items: int = 12, max_total_chars: int = 600) -> List[str]:
    """
    슬라이드 메타데이터에서 영상 화면(중앙)과 매칭에 사용할 텍스트를 추립니다.

    ★ v7: PPT 영역 기반 추출 (위치로 분류된 region_texts 사용)
    포함 영역:
      · center (중앙 콘텐츠) — 영상 화면과 1:1 매칭 대상
      · narration (하단 내레이션) — 강사가 말하는 핵심 키워드, 매칭 보조 신호
    제외 영역:
      · header (상단: 과정명/회차명/화면번호) — 메타정보
      · index (좌측 INDEX) — 모든 슬라이드에 같은 텍스트라 변별력 0
      · rdesc (우측 화면설명) — 메타정보 ("내용 제시" 등)

    region_texts가 없는 구버전 메타데이터(레거시)는 all_texts 사용으로 폴백.
    """
    out: List[str] = []
    total = 0
    seen = set()

    # 영역 기반 추출 (region_texts 우선)
    region_texts = info.get("region_texts") or {}
    if region_texts:
        # 중앙 텍스트가 가장 중요 → 먼저 추가
        text_pool: List[str] = []
        text_pool.extend(region_texts.get("center", []))
        text_pool.extend(region_texts.get("narration", []))
    else:
        # 폴백: 전체 텍스트 사용 (구버전 호환)
        text_pool = info.get("all_texts", []) or []

    screen_no = info.get("screen_no", "") or ""

    for raw in text_pool:
        if not raw:
            continue
        t = raw.strip()
        if not t or t == screen_no:
            continue

        # 보일러플레이트
        low = t.lower()
        if low in ("bgm", "자막영역 침범금지", "그림", "그림 1", "그림1",
                   "내용 제시", "내용을 입력하십시오"):
            continue
        if t.startswith("그림"):
            continue
        if t.startswith("#") and len(t) < 6:
            continue

        # 내레이션 텍스트 정제: "교수:" 표식 제거
        cleaned_text = re.sub(r'^\s*교수\s*:\s*', '', t, flags=re.MULTILINE)

        # 라인 단위로 풀어서 #N 토큰 제거
        cleaned_lines = []
        for line in cleaned_text.splitlines():
            ln = re.sub(r'^\s*#\d+\s*', '', line).strip()
            ln = re.sub(r'\s+#\d+\s*$', '', ln).strip()
            ln = re.sub(r'\s*#\d+\s*', ' ', ln).strip()
            if ln and len(ln) > 1 and len(ln) < 200:
                cleaned_lines.append(ln)
        for ln in cleaned_lines:
            if ln in seen:
                continue
            seen.add(ln)
            out.append(ln)
            total += len(ln)
            if len(out) >= max_items or total >= max_total_chars:
                return out
    return out


def _extract_center_texts(info: dict) -> List[str]:
    """
    슬라이드의 **중앙 영역 텍스트만** 반환합니다.
    화면 텍스트 검증(이슈 비교) 시, 메타 영역 텍스트가 false positive를
    유발하지 않도록 중앙 영역만 깨끗하게 추리는 용도.
    """
    region_texts = info.get("region_texts") or {}
    pool = region_texts.get("center", []) if region_texts else []
    out: List[str] = []
    seen = set()
    for raw in pool:
        if not raw:
            continue
        t = raw.strip()
        if not t:
            continue
        low = t.lower()
        if low in ("bgm", "자막영역 침범금지", "내용 제시", "내용을 입력하십시오"):
            continue
        if t.startswith("그림"):
            continue
        if t.startswith("#") and len(t) < 6:
            continue
        # #N 제거
        cleaned_lines = []
        for line in t.splitlines():
            ln = re.sub(r'\s*#\d+\s*', ' ', line).strip()
            if ln and len(ln) > 1 and len(ln) < 200:
                cleaned_lines.append(ln)
        for ln in cleaned_lines:
            if ln not in seen:
                seen.add(ln)
                out.append(ln)
    return out


def build_slide_display_meta(metadata: List[dict], meta_offset: int = 0) -> List[dict]:
    """
    스토리보드 매핑에서 '콘텐츠 슬라이드'만 추려 Vision API 프롬프트에 넣을
    간결한 display metadata 리스트를 반환합니다.

    meta_offset 만큼 앞 슬라이드(용어집/설계가이드 등)를 건너뜁니다.
    결과 리스트의 인덱스는 sb_images 리스트의 인덱스와 1:1 대응합니다.

    Returns
    -------
    List[dict] — 각 항목: {
        "slide_idx": int,          # sb_images에서의 0-based 인덱스
        "slide_no":  int,          # 사용자에게 보여줄 1-based 번호
        "screen_no": str,          # "03_02" 등
        "chapter":   str,          # 좌측상단 1줄 검증 기준
        "subtopic":  str,          # 좌측상단 2줄 검증 기준
        "title":     str,          # 슬라이드 타이틀
    }
    """
    display = []
    for i, info in enumerate(metadata[meta_offset:]):
        display.append({
            "slide_idx": i,
            "slide_no":  i + 1,
            "screen_no": info.get("screen_no", "") or "",
            "chapter":   info.get("chapter", "") or "",
            "subtopic":  info.get("subtopic", "") or info.get("title", "") or "",
            "title":     info.get("title", "") or "",
            # ★ v4: 텍스트 콘텐츠 기반 매칭용 (중앙 + 내레이션)
            "content_texts": _extract_content_texts(info),
            # ★ v7: 화면 텍스트 검증 전용 — 중앙 영역만 (내레이션 제외)
            "center_texts":  _extract_center_texts(info),
            "narration": (info.get("narration", "") or "")[:200],
        })
    return display


# ─────────────────────────────────────────────
# 화면 텍스트 맞춤법 검사
# ─────────────────────────────────────────────

_OCR_SPELL_PROMPT = f"""당신은 초정밀 한국어 OCR 및 맞춤법 교정 전문가입니다.
제공된 비디오 캡처 이미지들을 면밀히 분석하세요.

【작업 순서】
1. 각 이미지에서 보이는 **모든 한국어 텍스트**를 빠짐없이 `transcription`에 기록하세요.
   - 자막, 제목, 자료 화면, UI 라벨, 워터마크 등 위치 무관 전부 포함
   - 글자가 작거나 흐릿해도 최선을 다해 판독
2. 추출한 텍스트에서 맞춤법·오타 오류를 찾아 `corrections`에 담으세요.
3. **중복 제거**: 이전 이미지들과 **완전히 동일한 자막/텍스트**만 보인다면
   해당 이미지는 `corrections: []`로 비워두세요.

{_KO_SPELL_RULES}
{_SPELL_FEWSHOT}

【교정 규칙】
- 문장 전체를 반환하되, 오류 부분만 <red>단어</red>로 감싸세요.
- 확실한 오류만 교정 (불확실 → 원문 유지).
- OCR 판독 불확실한 글자는 `transcription`에는 추측, `corrections`에는 넣지 마세요.
- 이미지에 한국어 텍스트가 없거나 오류 없으면 `corrections: []`.
"""

_OCR_SCHEMA = {
    "name": "ocr_spell_results",
    "strict": True,
    "schema": {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "results": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "id": {"type": "integer"},
                        "matched_slide_number": {"type": "integer"},
                        "match_confidence": {"type": "string"},
                        "transcription": {"type": "string"},
                        "corrections": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "additionalProperties": False,
                                "properties": {
                                    "original":  {"type": "string"},
                                    "corrected": {"type": "string"},
                                    "reason":    {"type": "string"},
                                },
                                "required": ["original", "corrected", "reason"],
                            },
                        },
                    },
                    "required": ["id", "matched_slide_number", "match_confidence", "transcription", "corrections"],
                },
            }
        },
        "required": ["results"],
    },
}

# ─────────────────────────────────────────────
# ★ v3: 4항목 전면 검증 프롬프트 (스토리보드-영상 QA 스킬 정합)
#   (1) 좌측상단 2줄 (챕터/소주제) 텍스트 일치
#   (2) 우측상단 STEP 로고 존재
#   (3) 중앙 화면 텍스트 일치
#   (4) 한국어 맞춤법/띄어쓰기
# ─────────────────────────────────────────────
_SB_FULLMATCH_PROMPT = f"""당신은 교육 영상이 스토리보드와 일치하는지 검수하는 전문가입니다.
각 비디오 프레임에 대해 아래 4가지 항목을 모두 점검하고 발견된 이슈를 보고하세요.

【영상 프레임의 레이아웃 — 반드시 숙지】
교육 영상은 다음 3개 영역으로 구성됩니다:

  ┌────────────────────────────────┬──────────┐
  │ 좌측상단 2줄:                  │ 우측상단 │
  │   1행(위) = 챕터명             │  STEP    │
  │   2행(아래) = 소주제            │  로고    │
  ├────────────────────────────────┴──────────┤
  │                                            │
  │        중앙 화면 (본문 컨텐츠)             │
  │  (다이어그램, 도표, 설명 텍스트, 강사 상반신) │
  │                                            │
  └────────────────────────────────────────────┘

【⚠️ 중요 — 영상과 스토리보드 이미지의 비교 방식】
이 프롬프트에 첨부된 **스토리보드 슬라이드 이미지는 이미 "중앙 콘텐츠 노란 박스 영역"만 잘라낸 상태**
입니다. 좌측 INDEX, 상단 헤더, 우측 화면설명, 하단 내레이션 등은 모두 잘려있고
**영상 화면과 1:1 비교 가능한 영역**만 보입니다.

따라서:
  - 영상 프레임 중앙 ↔ 슬라이드 중앙 크롭 이미지를 직접 시각·텍스트 비교하세요.
  - **디자인·레이아웃·색상은 영상과 스토리보드가 다를 수 있습니다** (영상은 깔끔한
    콘텐츠 화면, 스토리보드는 디자인 시안). 그러나 **표현되는 텍스트·도표·그림의
    의미적 콘텐츠는 동일**해야 합니다.
  - 매칭 판단의 1순위는 **텍스트 콘텐츠 일치**, 2순위는 **시각적 구조 일치**입니다.

영상의 '좌측상단 2줄'(챕터/소주제) 은 스토리보드 이미지에는 없으므로 프롬프트로
별도 제공되는 [기준 챕터] / [기준 소주제] 텍스트와 비교합니다.

【⚠️ 중요 — 영상과 스토리보드의 디자인 차이】
영상의 중앙 화면과 스토리보드 슬라이드의 중앙 영역은 **디자인 구조가 완전히 다릅니다.**
- 영상: 깔끔한 콘텐츠 중심 화면(다이어그램·도표·텍스트가 큼직하게 배치)
- 스토리보드: 표 템플릿 안에 작은 썸네일 형태로 콘텐츠가 들어 있음
이 둘을 시각적 레이아웃·색상·구도로 비교하면 절대 일치하지 않습니다.
**오직 화면에 표시된 텍스트(단어·문장·도형 라벨) 내용**으로만 매칭을 판단하세요.

【슬라이드 매칭 방법 — 텍스트 콘텐츠 우선】
1. 영상 프레임 중앙에서 보이는 모든 텍스트(제목·라벨·본문 단어)를 빠짐없이 읽어내세요.
2. 프롬프트 상단에 슬라이드별로 제공된 **[중앙 영역 텍스트]** 목록과 우선 비교하세요.
   중앙 텍스트의 핵심 단어·구절이 영상 프레임에도 동일하게 보이는 슬라이드를 매칭합니다.
3. 중앙 텍스트만으로 판별이 어렵다면 **[내레이션]** 의 핵심 키워드도 보조 신호로 활용하세요.
   (강사가 말하는 내용이 화면에 텍스트·라벨로 등장하는 경우가 많음)
4. 디자인이 달라 보여도 텍스트 콘텐츠가 일치하면 매칭됩니다.
   예) 중앙 텍스트가 ["TSV 공정 개요", "실리콘 관통 전극", "임시 본딩"] 인데
       영상 프레임에 "TSV 공정 개요"라는 큰 제목과 "실리콘 관통 전극"이 보이면 → 매칭.
5. `matched_slide_number`: 1-based 슬라이드 번호. 텍스트가 어느 슬라이드와도
   유의미하게 겹치지 않으면 0 (매칭 없음).
6. `match_confidence`:
   - "exact"  ... 중앙 텍스트의 핵심 키워드 3개 이상이 영상에 동일하게 나타남
   - "high"   ... 핵심 키워드 2개 일치 + 다른 슬라이드보다 명확히 우위
   - "medium" ... 키워드 1~2개 일치, 다른 후보와 비슷하게 비등
   - "low"    ... 키워드 1개만 약하게 일치, 또는 일부 단어만 부분 일치
   - "none"   ... 어느 슬라이드와도 텍스트가 매칭되지 않음 (이때 matched_slide_number=0)

⚠️ **매우 중요 — 매칭 판단 가이드**
- "디자인이 다르니까 매칭 없음(0)"으로 처리하면 안 됩니다. 영상과 스토리보드는
  **원래** 디자인이 다른 게 정상입니다. 디자인 차이만으로 0을 답하지 마세요.
- "확신이 없으니까 0"도 잘못된 판단입니다. 핵심 키워드 1개라도 일치하면 'low'로
  매칭하고, 'none'(0)은 정말로 영상 텍스트가 어느 슬라이드와도 단 한 단어도
  겹치지 않을 때만 사용하세요.
- 영상에 보이는 한 단어가 어느 슬라이드의 [본문 텍스트] 항목 안에 들어 있는지
  하나씩 체크하세요. 부분 일치(예: "TSV"가 "TSV 공정 개요" 안에 포함)도 매칭입니다.
- 같은 키워드가 여러 슬라이드에 있으면 가장 많은 키워드가 겹치는 슬라이드를 고르세요.
- 영상 프레임이 강사만 보이고 콘텐츠가 적어 보여도, 자막·라벨·작은 텍스트라도
  매칭에 활용하세요. transcription을 비워두지 마세요.

【검사 4항목】
각 프레임에서 아래를 모두 점검하여 `issues` 배열에 기록. 이슈가 없으면 빈 배열.
각 이슈의 `issue_type` 은 반드시 다음 4개 중 하나:
  - "좌측상단"   ... 영상 좌측상단 2줄의 챕터/소주제 텍스트가 기준과 다름
  - "STEP 로고"  ... 영상 우측상단에 STEP 로고가 없거나 왜곡/잘림
  - "화면 텍스트" ... 중앙 화면 영역 텍스트가 스토리보드 중앙과 다름
  - "맞춤법"     ... 영상에 표시된 한국어 텍스트의 맞춤법/띄어쓰기 오류

(1) 좌측상단 2줄 [issue_type="좌측상단"]
   - 영상 좌측상단의 윗줄(챕터명) / 아랫줄(소주제)을 OCR로 읽어내세요.
   - 프롬프트 상단에 슬라이드별로 제공된 [기준 챕터] / [기준 소주제] 와 글자 단위로 비교.
   - 오타·띄어쓰기·대소문자 불일치 모두 이슈. `before`=영상에서 보이는 텍스트, `after`=기준 텍스트.
   - 영상에 해당 줄이 아예 없거나 잘린 경우에도 이슈.

(2) 우측상단 STEP 로고 [issue_type="STEP 로고"]
   - 영상 우측상단 구석에 파란색 계단 모양 아이콘 + "STEP" 글자 로고가 있어야 합니다.
   - 없거나 잘림/왜곡/다른 로고가 보이면 이슈.
   - `before`="(로고 없음)" 또는 실제 상태 설명, `after`="STEP 로고",
     `reason`="우측상단 STEP 로고가 누락/왜곡되었습니다." 형식.

(3) 중앙 화면 텍스트 [issue_type="화면 텍스트"]
   - 영상 프레임의 텍스트와 매칭된 슬라이드의 **[중앙 영역 텍스트]** 항목들만을
     글자 단위로 비교합니다. (좌측 INDEX·상단 헤더·우측 화면설명·하단 내레이션은
     영상에 나오지 않으므로 비교 대상이 아닙니다.)
   - **디자인·레이아웃·색상·도형 모양 차이는 절대 이슈가 아닙니다.**
     영상의 콘텐츠 화면과 스토리보드의 표 안 썸네일은 원래 디자인이 완전히 다릅니다.
     이런 디자인 차이는 보고하지 마세요.
   - 정말 보고할 것: **중앙 영역 텍스트 단어·문장의 글자 단위 불일치**만.
     예) 슬라이드 중앙 텍스트가 "구리 충전 기술"인데 영상엔 "구리충전기술"로 띄어쓰기 다름 → 이슈.
     예) 슬라이드 중앙 텍스트가 "TSV 공정"인데 영상엔 "TVS 공정"으로 오타 → 이슈.
   - 영상에 보이지만 중앙 영역 텍스트 목록에 없는 단어는 이슈로 보고하지 마세요
     (영상에만 있는 추가 라벨/주석일 수 있음).
   - matched_slide_number=0 (매칭 실패)일 때는 화면 텍스트 이슈를 보고하지 마세요.
     비교 기준이 없으므로 판단할 수 없습니다.
   - `before`=영상에서 보이는 텍스트, `after`=슬라이드 중앙의 올바른 텍스트.
   - `reason`="[스토리보드 N번 슬라이드와 불일치] ~" 형식으로 매칭 슬라이드 번호 명시.

(4) 맞춤법/띄어쓰기 [issue_type="맞춤법"]
   - 영상 중앙 화면에 보이는 한국어 텍스트의 맞춤법·띄어쓰기 오류.
   - 스토리보드에도 같은 오류가 있다면 `reason` 앞에 "[스토리보드 원본 오타] " 를,
     영상에만 있는 오류면 "[영상 텍스트 오류] " 를 붙이세요.
   - 한국어 기준: "첫번째"→"첫 번째", "구리충전기술"→"구리 충전 기술" 등 국립국어원 규정 적용.

{_KO_SPELL_RULES}

【공통 출력 규칙】
- `before` / `after` 는 문장 전체를 쓰되, 수정 부분만 <red>태그</red>로 감싸세요.
  (단, "STEP 로고" 이슈는 태그 없이 자유 문장)
- `transcription`: 프레임 중앙 화면에서 읽힌 모든 텍스트(좌측상단·로고 영역 제외).
- 다음은 **오류가 아닙니다** — 무시:
  · 강사 상반신, 3D 애니메이션 배경, 캐릭터·일러스트
  · 폰트/색상/위치만 다르고 텍스트 내용은 같은 경우
  · 스토리보드에만 있는 #0, #1 등 애니메이션 순서 표시
- 불확실한 이슈는 기록하지 말고 건너뛰세요. false positive보다 누락이 낫습니다.

【⚠️ 미완성 프레임(애니메이션 중간) 처리 — 매우 중요】
교육 영상에는 페이드인·슬라이드 인·도형 그리기·텍스트 한 글자씩 등장 등의
애니메이션이 흔합니다. 다음 징후가 보이면 그 프레임은 **"화면 텍스트" 이슈를
보고하지 말고** match_confidence 만 'low' 또는 'medium' 으로 매칭하세요:
  · 텍스트가 부분적으로만 보이거나 잘려있다
  · 다이어그램·도형·화살표가 일부만 그려져 있다
  · 페이드/투명도가 진행 중인 듯한 흐릿함
  · 박스나 글자가 화면 가장자리에서 들어오고 있는 모습
이런 미완성 상태와 스토리보드(완성 상태)를 비교한 차이는 영상 오류가 아니라
애니메이션 진행 중일 뿐입니다. **반드시 화면 텍스트 이슈로 보고하지 마세요.**
단, "좌측상단"·"STEP 로고"·"맞춤법"은 미완성 프레임이라도 정상적으로 점검합니다.
"""

_SB_FULLMATCH_SCHEMA = {
    "name": "sb_fullmatch_results",
    "strict": True,
    "schema": {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "results": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "id": {"type": "integer"},
                        "matched_slide_number": {"type": "integer"},
                        "match_confidence": {"type": "string"},
                        "transcription": {"type": "string"},
                        "issues": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "additionalProperties": False,
                                "properties": {
                                    "issue_type": {"type": "string"},
                                    "before":     {"type": "string"},
                                    "after":      {"type": "string"},
                                    "reason":     {"type": "string"},
                                },
                                "required": ["issue_type", "before", "after", "reason"],
                            },
                        },
                    },
                    "required": ["id", "matched_slide_number", "match_confidence", "transcription", "issues"],
                },
            }
        },
        "required": ["results"],
    },
}

_SB_MATCH_PROMPT = f"""당신은 영상 화면이 스토리보드 슬라이드와 1:1로 정확히 일치하는지 검수하는 전문가입니다.
스토리보드 슬라이드 이미지 일부가 참조용으로 먼저 제공되고, 이어서 비디오 캡처 이미지가 제공됩니다.

【스토리보드 템플릿 구조 — 반드시 숙지】
스토리보드 슬라이드는 표 형태 템플릿이며, 아래 5개 영역으로 나뉩니다:

  ┌──────────┬─────────────────────────┬──────────┐
  │  상단 헤더 (과정명, 회차명, 화면번호)             │
  ├──────────┼─────────────────────────┼──────────┤
  │          │                         │ 화면설명  │
  │  INDEX   │   ★ 중앙: 실제 영상 화면  │          │
  │ (목차)   │   (비교 대상은 여기만!)   │ 이미지   │
  │          │                         │  번호    │
  ├──────────┼─────────────────────────┼──────────┤
  │  내레이션 (교수: ...)                            │
  └──────────┴─────────────────────────┴──────────┘

  ❌ 영상에 없는 영역 (절대 비교하지 마세요):
     - 상단 헤더 (과정명, 회차명, 화면번호)
     - 좌측 INDEX (목차, 빨간 박스 강조 표시)
     - 우측 화면설명, 이미지 번호
     - 하단 내레이션 텍스트
  ✅ 영상에 있는 영역 (이것만 비교하세요):
     - 중앙 화면 이미지 영역의 텍스트와 그래픽

비디오 프레임에는 오직 중앙 화면 영역의 내용만 표시됩니다.
좌측 INDEX, 우측 화면설명, 상단 헤더, 하단 내레이션은 영상에 전혀 나타나지 않습니다.

【핵심 목표: 1:1 슬라이드 매칭】
각 비디오 프레임이 스토리보드의 몇 번 슬라이드에 해당하는지 정확히 판별하세요.

【작업 순서】
1. **슬라이드 매칭**: 비디오 이미지의 내용이 스토리보드 중앙 화면 영역과 일치하는 슬라이드를 찾기
   - `matched_slide_number`: 매칭되는 슬라이드 번호 (1부터 시작). 매칭 없으면 0.
   - `match_confidence`: 매칭 확신도
     · "exact" — 중앙 화면 영역의 레이아웃과 텍스트가 동일
     · "high" — 거의 동일 (사소한 차이만 있음)
     · "medium" — 유사하지만 텍스트 수정 또는 레이아웃 변경 있음
     · "low" — 매칭이 불확실함
     · "none" — 매칭되는 슬라이드 없음 (matched_slide_number=0일 때)
   - 중앙 화면의 레이아웃, 색상, 텍스트를 종합적으로 비교하세요.
2. **텍스트 비교**: 매칭된 슬라이드의 **중앙 화면 영역**과 비디오 이미지의 텍스트를 **글자 단위로** 비교
   - `transcription`: 비디오 이미지에서 읽히는 모든 한국어 텍스트를 기록
   - `corrections`: 중앙 화면 영역 텍스트와 다른 부분이 있으면 기재, 동일하면 빈 배열
   - ⚠️ 다시 강조: INDEX, 화면설명, 헤더, 내레이션의 텍스트는 절대 비교하지 마세요
3. 다음은 **오류가 아닙니다** — 무시:
   - 캐릭터, 배경, 일러스트 등 비텍스트 디자인 차이
   - 폰트 크기, 색상, 위치만 다른 경우 (텍스트 내용이 같으면 OK)
   - 애니메이션 효과로 인한 텍스트 일부 가림
   - #0, #1 등 애니메이션 순서 표시 (스토리보드에만 있고 영상에는 없음)

{_KO_SPELL_RULES}

【교정 규칙】
- original: 영상(비디오)에 어떻게 적혀 있는지 기재 (문장/구절 전체, 수정 부분만 <red>태그</red>)
- corrected: 스토리보드 중앙 화면 기준으로 어떻게 적혀야 하는지 기재 (수정 부분만 <red>태그</red>)
- reason: "[스토리보드 n번 슬라이드와 불일치] ~" 형식으로 사유 명기
"""

def spell_check_frames(
    frames, api_key, batch_size=2, model="gpt-5.4", max_workers: int = 4,
    storyboard_images: Optional[List[str]] = None,
    start_sb_idx: int = 0,
    verify_pass: bool = True,
    slide_metadata: Optional[List[dict]] = None,
    storyboard_match_images: Optional[List[str]] = None,   # ★ v10: 매칭 전용 (중앙 크롭)
):
    """
    프레임 이미지를 GPT Vision으로 OCR + 맞춤법 검사합니다.

    ★ 정확도 v2:
    - 히스토그램+SSIM 유사도로 슬라이드 후보를 사전 선별
    - verify_pass: True면 1차 결과를 2차 검증
    - 모든 프레임의 1:1 슬라이드 매칭 결과를 slide_map으로 반환

    ★ v3 (스토리보드-영상 QA 스킬 정합):
    - slide_metadata 제공 시 4항목 검증 (좌측상단/STEP 로고/화면 텍스트/맞춤법) 수행
      · 각 항목은 build_slide_display_meta() 결과와 동일 포맷:
        {"slide_idx": int, "slide_no": int, "screen_no": str,
         "chapter": str, "subtopic": str, "title": str}
      · 인덱스는 storyboard_images와 1:1 대응해야 함
    - slide_metadata=None 이면 기존 중앙-텍스트 비교 모드로 동작

    Returns
    -------
    (corrections, slide_map, final_sb_idx)
    - corrections : List[dict] — 텍스트 교정 결과
    - slide_map   : List[dict] — 모든 프레임의 슬라이드 매칭 정보
        각 항목: {
            "frame_time": "[01:23]",
            "matched_slide": 5,        # 0 = 매칭 없음
            "match_confidence": "exact",  # exact/high/medium/low/none
            "similarity_score": 0.85,  # 히스토그램+SSIM 점수
            "transcription": "프레임 텍스트",
            "has_corrections": True/False,
            "image_b64": "...",
        }
    - final_sb_idx : int — 마지막 매칭된 슬라이드 인덱스
    """
    if not frames:
        return [], [], start_sb_idx

    client = OpenAI(api_key=api_key)

    batches: List[List[dict]] = [
        frames[i:i + batch_size] for i in range(0, len(frames), batch_size)
    ]

    is_sb = bool(storyboard_images)
    use_fullmatch = bool(slide_metadata) and is_sb    # ★ v3: 4항목 전면 검증 모드

    if is_sb:
        max_workers = 1

    # 전체 슬라이드 매핑 결과를 수집할 리스트
    all_slide_map: List[dict] = []

    def _process_batch(batch_idx: int, batch: List[dict], current_sb_idx: int = 0) -> Tuple[List[dict], List[dict], int]:
        """Returns (corrections, batch_slide_map, next_sb_idx)"""
        if use_fullmatch:
            prompt_text = _SB_FULLMATCH_PROMPT
        elif is_sb:
            prompt_text = _SB_MATCH_PROMPT
        else:
            prompt_text = _OCR_SPELL_PROMPT
        content_items = [{"type": "text", "text": prompt_text}]
        next_sb_idx = current_sb_idx

        # 각 프레임의 사전 유사도 점수를 미리 계산
        frame_similarity_scores: Dict[int, Tuple[int, float]] = {}  # frame_idx -> (best_slide_idx, score)

        if is_sb:
            content_items.append({"type": "text", "text": "=== [탐색 대상 스토리보드 슬라이드 범위] ==="})
            # ★ v4: 영상-스토리보드 디자인 차이가 크므로 시각 유사도 사전 필터를 완화.
            #   후보 수를 늘리고, 부족하면 search_range 전체를 후보로 사용.
            # ★ FULLSCAN v1: 시작 슬라이드가 틀리거나 메타 제외가 어긋나도 놓치지 않도록
            # 모든 배치에서 처음부터 끝까지 전체 슬라이드를 스캔합니다.
            # GPT에는 전체를 모두 보내지 않고, 아래 시각/텍스트 사전 매칭으로 뽑힌 후보만 보냅니다.
            search_range = (0, len(storyboard_images))

            # ★ v10: 시각 매칭에는 중앙 크롭 이미지를 사용 (있는 경우)
            # 영상 화면과 1:1 디자인 일치하는 영역만 비교 → 매칭 정확도 크게 상승
            match_images = storyboard_match_images if storyboard_match_images else storyboard_images

            candidate_indices = set()
            for fi, frame_item in enumerate(batch):
                # ★ FULLSCAN v1: 원본 프레임 전체가 아니라 영상 중앙 콘텐츠 영역만 잘라서
                # 스토리보드 중앙 노란 박스 크롭과 비교합니다.
                frame_match_b64 = crop_video_content_region_b64(frame_item["base64"])
                top_matches = find_best_matching_slides(
                    frame_match_b64,
                    match_images,
                    search_range=search_range,
                    top_k=20,
                )
                for slide_idx, score in top_matches:
                    candidate_indices.add(slide_idx)
                if top_matches:
                    frame_similarity_scores[fi] = (top_matches[0][0], top_matches[0][1])

            # 후보가 부족하면 search_range 전체를 후보로
            min_required = max(6, min(10, search_range[1] - search_range[0]))
            if len(candidate_indices) < min_required:
                for i in range(search_range[0], min(search_range[1], search_range[0] + min_required)):
                    candidate_indices.add(i)

            sorted_candidates = sorted(candidate_indices)
            for sb_idx in sorted_candidates:
                if 0 <= sb_idx < len(storyboard_images):
                    # ★ v10: GPT 프롬프트에는 중앙 크롭 이미지를 보냄 (있으면)
                    # 영상과 1:1 비교 대상 영역만 보여주면 GPT가 디자인 차이로 헷갈리지 않음
                    sb_img = match_images[sb_idx] if sb_idx < len(match_images) else storyboard_images[sb_idx]

                    if use_fullmatch and 0 <= sb_idx < len(slide_metadata):
                        meta = slide_metadata[sb_idx]
                        center_texts = meta.get("center_texts", []) or []
                        narration_text = (meta.get("narration", "") or "").strip()

                        if center_texts:
                            center_lines = "\n".join(f"    - {t}" for t in center_texts[:10])
                        else:
                            center_lines = "    (중앙 영역 텍스트 추출 실패 — 슬라이드 이미지 참조)"

                        if narration_text:
                            narr_str = narration_text[:300] + ("..." if len(narration_text) > 300 else "")
                        else:
                            narr_str = "(내레이션 없음)"

                        hdr = (
                            f"[스토리보드 {sb_idx + 1}번 슬라이드 — "
                            f"화면번호 {meta.get('screen_no', '-') or '-'}]\n"
                            f"  [기준 챕터]   {meta.get('chapter',  '') or '(없음)'}\n"
                            f"  [기준 소주제] {meta.get('subtopic', '') or '(없음)'}\n"
                            f"  [중앙 영역 텍스트] ★ 영상 화면(중앙)과 1:1 비교 대상 ★\n"
                            f"{center_lines}\n"
                            f"  [내레이션] (참고용 — 영상 음성과 비교, 매칭에도 보조 활용)\n"
                            f"    {narr_str}\n"
                            f"  [슬라이드 이미지] ↓ 아래 이미지는 이 슬라이드의 **중앙 콘텐츠 노란 박스만 잘라낸** 모습입니다."
                        )
                        content_items.append({"type": "text", "text": hdr})
                    else:
                        content_items.append({"type": "text", "text": f"[스토리보드 {sb_idx + 1}번 슬라이드 — 중앙 영역]"})
                    content_items.append({
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{sb_img}", "detail": "high"},
                    })
            content_items.append({"type": "text", "text": "=== [비디오 프레임 이미지] ==="})

        for idx, f in enumerate(batch):
            content_items.append({"type": "text", "text": f"[비디오 {idx}번 이미지 — {f['time_str']}]"})
            content_items.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{f['base64']}", "detail": "high"},
            })

        def _call():
            _schema = _SB_FULLMATCH_SCHEMA if use_fullmatch else _OCR_SCHEMA
            return client.chat.completions.create(
                model=model,
                messages=[{"role": "user", "content": content_items}],
                temperature=0.0,
                response_format={"type": "json_schema", "json_schema": _schema},
            )

        corrections: List[dict] = []
        batch_slide_map: List[dict] = []

        # 기본 매핑 엔트리 (GPT 응답 없어도 프레임 기록은 남김)
        for idx, f in enumerate(batch):
            sim_info = frame_similarity_scores.get(idx)
            batch_slide_map.append({
                "frame_time": f["time_str"],
                "frame_time_sec": float(f.get("time", 0) or 0),   # ★ v8: 시간 보간용 숫자 시간
                "matched_slide": 0,
                "match_confidence": "none",
                "similarity_score": sim_info[1] if sim_info else 0.0,
                "transcription": "",
                "has_corrections": False,
                "image_b64": f["base64"],
            })

        try:
            response = call_with_retry(_call)
            parsed = json.loads(response.choices[0].message.content)
            for item in parsed.get("results", []):
                f_idx = item.get("id")
                matched_sb = item.get("matched_slide_number") or 0
                match_conf = item.get("match_confidence", "none") or "none"
                transcription = item.get("transcription", "") or ""

                if matched_sb and isinstance(matched_sb, int) and matched_sb > 0:
                    next_sb_idx = matched_sb - 1

                if f_idx is None or not (0 <= f_idx < len(batch)):
                    continue

                # 슬라이드 매핑 정보 업데이트
                batch_slide_map[f_idx]["matched_slide"] = matched_sb
                batch_slide_map[f_idx]["match_confidence"] = match_conf
                batch_slide_map[f_idx]["transcription"] = transcription

                # ★ v3 fullmatch 모드: issues 배열 파싱 (4개 issue_type 분류)
                # 구 모드: corrections 배열 파싱 (모두 "화면 텍스트"로 분류)
                frame_has_corrections = False
                issue_list = item.get("issues") if use_fullmatch else item.get("corrections", [])
                issue_list = issue_list or []

                for corr_item in issue_list:
                    orig = (corr_item.get("before") if use_fullmatch else corr_item.get("original")) or ""
                    corr = (corr_item.get("after")  if use_fullmatch else corr_item.get("corrected")) or ""
                    orig = orig.strip(); corr = corr.strip()
                    if not (orig and corr) or orig == corr:
                        # STEP 로고 이슈처럼 "원본과 기준이 같지만 상태가 다른" 경우를 놓치지 않도록,
                        # issue_type == "STEP 로고" 는 공란이 아니면 통과시킴
                        if use_fullmatch and (corr_item.get("issue_type") == "STEP 로고") and (orig or corr):
                            pass
                        else:
                            continue

                    if use_fullmatch:
                        it = (corr_item.get("issue_type") or "").strip()
                        if it not in ("좌측상단", "STEP 로고", "화면 텍스트", "맞춤법"):
                            it = "화면 텍스트"
                    else:
                        it = "화면 텍스트"

                    # ★ v4: 매칭 실패(matched_sb=0) 상태에서 "화면 텍스트" 이슈는
                    # 비교 기준이 없으므로 신뢰할 수 없음 → 드롭. (좌측상단/STEP/맞춤법은 슬라이드 매칭 없이도 검증 가능)
                    if it == "화면 텍스트" and (not matched_sb or matched_sb <= 0):
                        continue

                    frame_has_corrections = True
                    corrections.append({
                        "구분": it,
                        "시간": batch[f_idx]["time_str"],
                        "수정 전": _red_to_html(orig),
                        "수정 후": _red_to_html(corr),
                        "교정 사유": corr_item.get("reason") or "",
                        "image_b64": batch[f_idx]["base64"],
                        "_matched_sb": matched_sb,
                        "_issue_type": it,
                    })
                batch_slide_map[f_idx]["has_corrections"] = frame_has_corrections

        except Exception as e:
            print(f"화면 맞춤법 검사 오류 (배치 {batch_idx + 1}): {e}")
            try:
                corrections.extend(_fallback_text_mode_ocr(client, model, batch, batch_idx))
            except Exception as e2:
                print(f"  폴백도 실패: {e2}")

        return corrections, batch_slide_map, next_sb_idx

    # 실행 (스토리보드 모드: 순차)
    results: List[dict] = []
    slide_map_all: List[dict] = []
    final_sb_idx = start_sb_idx

    if max_workers <= 1 or len(batches) == 1:
        current_sb = start_sb_idx
        for i, b in enumerate(batches):
            corrs, smap, current_sb = _process_batch(i, b, current_sb)
            results.extend(corrs)
            slide_map_all.extend(smap)
        final_sb_idx = current_sb
    else:
        with ThreadPoolExecutor(max_workers=max_workers) as ex:
            futures = {ex.submit(_process_batch, i, b, 0): i for i, b in enumerate(batches)}
            indexed_corrs: Dict[int, List[dict]] = {}
            indexed_smap: Dict[int, List[dict]] = {}
            for fut in as_completed(futures):
                i = futures[fut]
                try:
                    corrs, smap, _ = fut.result()
                    indexed_corrs[i] = corrs
                    indexed_smap[i] = smap
                except Exception as e:
                    print(f"병렬 배치 {i + 1} 실패: {e}")
                    indexed_corrs[i] = []
                    indexed_smap[i] = []
            for i in sorted(indexed_corrs.keys()):
                results.extend(indexed_corrs[i])
                slide_map_all.extend(indexed_smap.get(i, []))

    # 중복 교정 제거
    seen = set()
    deduped = []
    for r in results:
        key = (
            re.sub(r"<[^>]+>", "", r["수정 전"]).strip(),
            re.sub(r"<[^>]+>", "", r["수정 후"]).strip(),
        )
        if key in seen:
            continue
        seen.add(key)
        deduped.append(r)

    # ★ v5/v6: 텍스트 콘텐츠 기반 폴백 매칭
    # GPT가 시각 차이를 이유로 매칭=0으로 답한 프레임을 transcription과
    # PPT 본문 텍스트의 직접 비교로 재시도. 디자인이 달라도 텍스트만 일치하면 매칭됨.
    # ★ v6: 내레이션을 매칭 풀에 포함, 임계값 0.30→0.20 완화, 디버깅용 top-5 점수 보존
    if slide_metadata and slide_map_all:
        unmatched_count = 0
        recovered_count = 0
        for entry in slide_map_all:
            transcription = entry.get("transcription", "") or ""
            # 매칭 안 된 프레임뿐 아니라, 매칭됐지만 confidence가 낮은 프레임도 점수 분포 저장
            if not transcription.strip():
                continue   # OCR 실패 — 매칭 불가
            sb_num, score, conf, top5 = text_based_match(transcription, slide_metadata)
            # 디버깅용 점수 분포 저장 (모든 프레임에 대해)
            entry["text_match_top5"] = top5
            entry["text_match_best_score"] = round(float(score), 3)

            # 매칭 실패한 프레임만 폴백 적용
            if entry.get("matched_slide", 0):
                continue
            unmatched_count += 1
            if sb_num > 0:
                entry["matched_slide"]    = sb_num
                entry["match_confidence"] = conf
                entry["similarity_score"] = round(float(score), 3)
                entry["text_matched"]     = True   # 디버깅용 표식
                recovered_count += 1
        if unmatched_count > 0:
            print(f"[v6 텍스트 폴백 매칭] 미매칭 {unmatched_count}개 중 {recovered_count}개 복구.")

    # ★ FIX v11: OCR/transcription이 비어 텍스트 폴백을 못 탄 프레임은
    # 사전 시각 유사도(best visual match)로 한 번 더 복구합니다.
    # 기존에는 transcription이 빈 경우 continue 되어 matched_slide=0으로 남았습니다.
    if slide_map_all and is_sb:
        visual_recovered = 0
        visual_sources = storyboard_match_images if storyboard_match_images else storyboard_images
        for entry in slide_map_all:
            if entry.get("matched_slide", 0):
                continue
            frame_b64 = entry.get("image_b64", "")
            if not frame_b64 or not visual_sources:
                continue
            # 전체 슬라이드에서 top1을 다시 계산. 영상도 중앙 콘텐츠 영역으로 크롭해서 비교합니다.
            frame_match_b64 = crop_video_content_region_b64(frame_b64)
            matches = find_best_matching_slides(frame_match_b64, visual_sources, search_range=None, top_k=1)
            if not matches and storyboard_images and visual_sources is not storyboard_images:
                matches = find_best_matching_slides(frame_b64, storyboard_images, search_range=None, top_k=1)
            if not matches:
                continue
            best_idx, best_score = matches[0]
            # ★ FULLSCAN v1: 매칭 없음으로 남기는 것보다 최고 후보를 low confidence로 남겨
            # 엑셀/화면에서 사람이 확인할 수 있게 합니다.
            entry["matched_slide"] = best_idx + 1
            entry["match_confidence"] = "low" if best_score < 0.45 else "medium"
            entry["similarity_score"] = round(float(best_score), 3)
            entry["visual_matched"] = True
            entry["forced_best_match"] = best_score < 0.30
            visual_recovered += 1
        if visual_recovered > 0:
            print(f"[v11 시각 폴백 매칭] OCR 미매칭 {visual_recovered}개 프레임 복구.")

    # ★ MOTION/ANCHOR v2: 스토리보드 순서 기반 앵커 보정
    # 구분/타이틀 화면이 잡히면 그 지점을 기준으로 이후 프레임을 주변 슬라이드에서 우선 재검색합니다.
    if slide_map_all and is_sb:
        seq_rec = _apply_storyboard_sequence_anchor(
            slide_map_all,
            slide_metadata,
            storyboard_images=storyboard_images,
            storyboard_match_images=storyboard_match_images,
        )
        if seq_rec > 0:
            print(f"[시퀀스 앵커 보정] {seq_rec}개 프레임을 스토리보드 진행 순서 기준으로 보정.")

    # ★ v8: 시간 기반 보간 매칭
    # 화면 중앙이 비어있거나 강사만 비추는 프레임은 텍스트로 매칭 불가능.
    # 영상은 시간 순서대로 슬라이드를 진행하므로, 앞뒤 프레임의 매칭 결과로 보간 가능.
    # 보수적 규칙: 직전·직후 매칭이 같은 슬라이드일 때만 보간 (둘이 다르면 전환 시점이라 모호).
    if slide_map_all and is_sb:
        rec = _interpolate_unmatched_by_time(slide_map_all, max_gap_seconds=12.0)
        if rec > 0:
            print(f"[v8 시간 보간 매칭] {rec}개 프레임을 앞뒤 매칭 결과로 보간.")

    # 2-pass 검증
    if verify_pass and is_sb and deduped:
        deduped = _verify_screen_corrections(client, model, deduped, storyboard_images)

    return deduped, slide_map_all, final_sb_idx


def _interpolate_unmatched_by_time(
    slide_map_all: List[dict],
    max_gap_seconds: float = 12.0,
) -> int:
    """
    매칭 안 된 프레임을 앞뒤로 인접한 매칭된 프레임의 결과로 보간합니다.

    ★ 보수적 규칙:
      · 직전(좌)과 직후(우) 매칭된 프레임이 **같은 슬라이드**일 때만 보간
      · 좌우 매칭 프레임 사이의 시간 간격이 max_gap_seconds 이내일 때만
        (간격이 너무 넓으면 다른 슬라이드 가능성이 커서 위험)

    이렇게 하면 슬라이드 전환 직전/직후의 빈 프레임이 안전하게 채워지고,
    챕터 전환 시점 같은 모호한 프레임은 매칭 실패 상태로 남깁니다.

    Returns
    -------
    int — 보간으로 매칭된 프레임 개수
    """
    if not slide_map_all:
        return 0

    n = len(slide_map_all)
    recovered = 0

    # 각 미매칭 프레임에 대해 좌우로 가장 가까운 매칭 찾기
    for i in range(n):
        if slide_map_all[i].get("matched_slide", 0):
            continue

        # 좌측 가장 가까운 매칭 프레임
        left_idx = -1
        for j in range(i - 1, -1, -1):
            if slide_map_all[j].get("matched_slide", 0):
                left_idx = j
                break

        # 우측 가장 가까운 매칭 프레임
        right_idx = -1
        for j in range(i + 1, n):
            if slide_map_all[j].get("matched_slide", 0):
                right_idx = j
                break

        if left_idx < 0 or right_idx < 0:
            continue   # 한쪽이라도 매칭이 없으면 보간 불가

        left_sb = slide_map_all[left_idx].get("matched_slide", 0)
        right_sb = slide_map_all[right_idx].get("matched_slide", 0)

        if left_sb != right_sb:
            continue   # 양쪽 슬라이드 번호가 다르면 보간 안 함 (전환 시점)

        # 시간 거리 체크
        try:
            left_time = float(slide_map_all[left_idx].get("frame_time_sec") or 0)
            right_time = float(slide_map_all[right_idx].get("frame_time_sec") or 0)
            if right_time - left_time > max_gap_seconds:
                continue
        except Exception:
            pass

        # 보간 적용
        slide_map_all[i]["matched_slide"]    = left_sb
        slide_map_all[i]["match_confidence"] = "low"
        slide_map_all[i]["similarity_score"] = max(
            slide_map_all[i].get("similarity_score", 0) or 0, 0.30
        )
        slide_map_all[i]["interpolated"]     = True   # 시간 보간 표식
        recovered += 1

    return recovered


def _verify_screen_corrections(
    client: OpenAI,
    model: str,
    corrections: List[dict],
    storyboard_images: List[str],
) -> List[dict]:
    """
    ★ 정확도 v2: 2-pass 검증
    1차에서 발견된 교정 결과를 재검증합니다.
    각 교정 항목에 대해 해당 프레임 이미지와 매칭된 슬라이드를 다시 비교하여
    실제로 오류인지 확인합니다. False positive를 제거합니다.
    """
    if not corrections:
        return corrections

    # 검증이 필요한 항목만 선별 (이미지와 매칭 슬라이드가 있는 것)
    # ★ v3: 좌측상단/STEP 로고는 스토리보드 중앙과 무관하므로 검증 대상에서 제외 (auto_pass)
    to_verify = []
    auto_pass = []
    for corr in corrections:
        it = corr.get("_issue_type") or corr.get("구분") or ""
        if it in ("좌측상단", "STEP 로고"):
            auto_pass.append(corr)
            continue
        sb_num = corr.get("_matched_sb")
        has_image = bool(corr.get("image_b64"))
        if sb_num and isinstance(sb_num, int) and sb_num > 0 and has_image:
            sb_idx = sb_num - 1
            if 0 <= sb_idx < len(storyboard_images):
                to_verify.append((corr, sb_idx))
                continue
        auto_pass.append(corr)

    if not to_verify:
        return corrections

    # 배치로 검증 (최대 5개씩)
    VERIFY_BATCH = 5
    verified = list(auto_pass)

    verify_prompt = """당신은 영상-스토리보드 비교 검증 전문가입니다.
이전 검수에서 발견된 차이점이 실제 오류인지 재확인해 주세요.

【영상 레이아웃】
- 좌측상단 2줄: 챕터명(위) / 소주제(아래)
- 우측상단: STEP 로고 (파란 계단 모양 + "STEP" 글자)
- 중앙: 본문 컨텐츠 (강사·3D배경 포함 가능)

【스토리보드 참고】스토리보드 슬라이드에는 좌측 INDEX·우측 화면설명·상단 헤더·하단 내레이션
등 영상에 나타나지 않는 영역이 있습니다. 슬라이드의 중앙 화면 영역만 영상 중앙과 대응합니다.
#0, #1 등 애니메이션 순서 표시는 영상에는 없으므로 오류로 판단하지 마세요.

【issue_type 별 검증 포인트】
- "좌측상단": 영상 좌측상단 2줄 텍스트와 "수정 후"(기준 챕터/소주제)가 실제로 다른지 확인.
- "STEP 로고": 영상 우측상단에 STEP 로고가 실제로 누락/왜곡되었는지 확인.
- "화면 텍스트": 영상 중앙과 스토리보드 중앙의 텍스트가 실제로 다른지 확인.
- "맞춤법": 지적된 부분이 실제 맞춤법/띄어쓰기 오류인지 확인.

【⚠️ 미완성 프레임 가드】
영상 프레임이 애니메이션 진행 중(텍스트가 부분적으로만 등장, 도형이 일부만 그려짐,
페이드 진행 중)이면 스토리보드(완성 상태)와 차이가 나는 게 당연합니다.
이런 차이는 영상의 오류가 아니므로 "화면 텍스트" 이슈는 "confirmed": false 로 처리하세요.
("좌측상단"·"STEP 로고"·"맞춤법"은 정상적으로 검증.)

각 항목에 대해 실제 오류면 "confirmed": true, 오류가 아니면 "confirmed": false.
"""

    verify_schema = {
        "name": "verify_results",
        "strict": True,
        "schema": {
            "type": "object",
            "additionalProperties": False,
            "properties": {
                "verifications": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "additionalProperties": False,
                        "properties": {
                            "index":     {"type": "integer"},
                            "confirmed": {"type": "boolean"},
                            "note":      {"type": "string"},
                        },
                        "required": ["index", "confirmed", "note"],
                    },
                }
            },
            "required": ["verifications"],
        },
    }

    for batch_start in range(0, len(to_verify), VERIFY_BATCH):
        batch = to_verify[batch_start:batch_start + VERIFY_BATCH]

        content_items = [{"type": "text", "text": verify_prompt}]
        for vi, (corr, sb_idx) in enumerate(batch):
            # 스토리보드 슬라이드
            content_items.append({"type": "text", "text": f"--- 항목 {vi} ---"})
            content_items.append({"type": "text", "text": f"[스토리보드 {sb_idx + 1}번 슬라이드]"})
            content_items.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{storyboard_images[sb_idx]}", "detail": "high"},
            })
            # 비디오 프레임
            content_items.append({"type": "text", "text": f"[비디오 프레임 — {corr['시간']}]"})
            content_items.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{corr['image_b64']}", "detail": "high"},
            })
            plain_before = re.sub(r"<[^>]+>", "", corr["수정 전"])
            plain_after = re.sub(r"<[^>]+>", "", corr["수정 후"])
            content_items.append({
                "type": "text",
                "text": f"제시된 차이: \"{plain_before}\" → \"{plain_after}\"\n사유: {corr['교정 사유']}",
            })

        try:
            def _verify_call():
                return client.chat.completions.create(
                    model=model,
                    messages=[{"role": "user", "content": content_items}],
                    temperature=0.0,
                    response_format={"type": "json_schema", "json_schema": verify_schema},
                )

            resp = call_with_retry(_verify_call)
            parsed = json.loads(resp.choices[0].message.content)
            confirmations = {v["index"]: v["confirmed"] for v in parsed.get("verifications", [])}

            for vi, (corr, _) in enumerate(batch):
                if confirmations.get(vi, True):  # 기본값 True (확인 안 되면 유지)
                    verified.append(corr)
                else:
                    print(f"  ✂️ 2차 검증에서 제거: {re.sub(r'<[^>]+>', '', corr['수정 전'])[:40]}")
        except Exception as e:
            print(f"2차 검증 오류 (유지): {e}")
            for corr, _ in batch:
                verified.append(corr)

    return verified


def _fallback_text_mode_ocr(client, model, batch, batch_idx):
    """OCR 텍스트 모드 폴백."""
    content_items = [{"type": "text", "text": _OCR_SPELL_PROMPT + (
        "\n응답은 순수 JSON 배열만 반환:\n"
        '[{"id":0,"transcription":"…","corrections":[{"original":"…","corrected":"…","reason":"…"}]}]'
    )}]
    for idx, f in enumerate(batch):
        content_items.append({"type": "text", "text": f"[{idx}번 이미지 — {f['time_str']}]"})
        content_items.append({
            "type": "image_url",
            "image_url": {"url": f"data:image/jpeg;base64,{f['base64']}", "detail": "high"},
        })

    response = client.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": content_items}],
        temperature=0.0,
    )
    content = _strip_json_fences(response.choices[0].message.content.strip())
    parsed = json.loads(content)
    if isinstance(parsed, dict):
        parsed = parsed.get("results", next(iter(parsed.values()), []))

    out = []
    for item in parsed:
        f_idx = item.get("id")
        if f_idx is None or not (0 <= f_idx < len(batch)):
            continue
        for corr_item in item.get("corrections", []):
            orig = (corr_item.get("original") or "").strip()
            corr = (corr_item.get("corrected") or "").strip()
            if orig and corr and orig != corr:
                out.append({
                    "구분": "화면 텍스트",
                    "시간": batch[f_idx]["time_str"],
                    "수정 전": _red_to_html(orig),
                    "수정 후": _red_to_html(corr),
                    "교정 사유": corr_item.get("reason", ""),
                    "image_b64": batch[f_idx]["base64"],
                })
    return out


# ─────────────────────────────────────────────
# 타임스탬프 → 프레임 추출
# ─────────────────────────────────────────────

def extract_frame_at_timestamp(video_path, time_str):
    """단일 타임스탬프의 프레임을 Base64로 반환."""
    try:
        ts = time_str.strip("[]")
        parts = list(map(int, ts.split(":")))
        if len(parts) == 2:
            seconds = parts[0] * 60 + parts[1]
        else:
            seconds = parts[0] * 3600 + parts[1] * 60 + parts[2]

        cap = cv2.VideoCapture(video_path)
        fps = cap.get(cv2.CAP_PROP_FPS) or 30.0
        cap.set(cv2.CAP_PROP_POS_FRAMES, int(seconds * fps))
        ret, frame = cap.read()
        cap.release()
        if ret:
            return encode_image(frame)
    except Exception as e:
        print(f"프레임 추출 오류 ({time_str}): {e}")
    return ""


def attach_frames_to_audio_results(audio_results, video_path):
    """
    음성 교정 결과에 영상 프레임을 image_b64로 추가합니다.
    (VideoCapture를 한 번만 열고 여러 타임스탬프를 읽어 I/O 최소화)
    """
    if not audio_results:
        return audio_results

    cap = cv2.VideoCapture(video_path)
    fps = cap.get(cv2.CAP_PROP_FPS) or 30.0

    def _ts_to_sec(time_str: str) -> Optional[float]:
        try:
            parts = list(map(int, time_str.strip("[]").split(":")))
            if len(parts) == 2:
                return parts[0] * 60 + parts[1]
            return parts[0] * 3600 + parts[1] * 60 + parts[2]
        except Exception:
            return None

    for r in audio_results:
        if r.get("image_b64"):
            continue
        seconds = _ts_to_sec(r.get("시간", ""))
        if seconds is None:
            continue
        try:
            cap.set(cv2.CAP_PROP_POS_FRAMES, int(seconds * fps))
            ret, frame = cap.read()
            if ret:
                r["image_b64"] = encode_image(frame)
            else:
                r["image_b64"] = ""
        except Exception as e:
            print(f"프레임 추출 오류 ({r.get('시간')}): {e}")
            r["image_b64"] = ""

    cap.release()
    return audio_results


# ─────────────────────────────────────────────
# ★ 정확도 v2: 슬라이드 커버리지 리포트
# ─────────────────────────────────────────────


def _apply_storyboard_sequence_anchor(
    slide_map_all: List[dict],
    slide_metadata: Optional[List[dict]],
    storyboard_images: Optional[List[str]] = None,
    storyboard_match_images: Optional[List[str]] = None,
    *,
    local_back: int = 3,
    local_forward: int = 18,
) -> int:
    """
    스토리보드 진행 순서를 이용해 흔들리는 매칭을 안정화합니다.

    원칙:
    1. 영상은 대체로 스토리보드 순서대로 진행되므로, 직전 확정 슬라이드 근처 후보를 우선합니다.
    2. 섹션 타이틀/구분 화면이 잡히면 그 슬라이드를 앵커로 삼습니다.
    3. 낮은 확신도의 갑작스러운 역주행/과도한 점프는 직전 앵커 주변에서 다시 이미지 매칭합니다.
    """
    if not slide_map_all:
        return 0
    visual_sources = storyboard_match_images or storyboard_images or []
    if not visual_sources:
        return 0

    changed = 0
    last_good: Optional[int] = None
    confidence_rank = {"none": 0, "low": 1, "medium": 2, "high": 3, "exact": 4}

    for entry in slide_map_all:
        cur = int(entry.get("matched_slide") or 0)
        conf = entry.get("match_confidence", "none") or "none"
        rank = confidence_rank.get(conf, 0)

        if cur > 0 and rank >= 3:
            last_good = cur
            entry["sequence_anchor"] = True
            continue

        if last_good is None:
            if cur > 0:
                last_good = cur
            continue

        suspicious = False
        if cur <= 0:
            suspicious = True
        elif rank <= 2 and (cur < last_good - local_back or cur > last_good + local_forward):
            suspicious = True

        if suspicious:
            frame_b64 = entry.get("image_b64", "")
            if frame_b64:
                start = max(0, last_good - 1 - local_back)
                end = min(len(visual_sources), last_good - 1 + local_forward + 1)
                frame_match_b64 = crop_video_content_region_b64(frame_b64)
                local_matches = find_best_matching_slides(
                    frame_match_b64,
                    visual_sources,
                    search_range=(start, end),
                    top_k=1,
                )
                if local_matches:
                    best_idx, best_score = local_matches[0]
                    new_slide = best_idx + 1
                    if new_slide != cur:
                        entry["matched_slide_before_sequence"] = cur
                        entry["matched_slide"] = new_slide
                        entry["match_confidence"] = "low" if best_score < 0.45 else "medium"
                        entry["similarity_score"] = round(float(best_score), 3)
                        entry["sequence_corrected"] = True
                        changed += 1
                        cur = new_slide

        if cur > 0 and cur >= last_good - local_back:
            last_good = max(last_good, cur)

    return changed

def compute_slide_coverage(
    slide_map: List[dict],
    total_slides: int,
) -> dict:
    """
    slide_map으로부터 스토리보드 슬라이드 커버리지 리포트를 생성합니다.

    Returns
    -------
    {
        "total_slides": 전체 슬라이드 수,
        "matched_slides": 영상에서 확인된 슬라이드 번호 set,
        "unmatched_slides": 영상에서 확인되지 않은 슬라이드 번호 set,
        "coverage_pct": 커버리지 비율 (0~100),
        "unmatched_frames": 슬라이드에 매칭되지 않은 프레임 수,
        "total_frames": 전체 프레임 수,
        "slide_frame_count": {슬라이드번호: 해당 슬라이드에 매칭된 프레임 수},
        "confidence_summary": {확신도: 프레임 수},
    }
    """
    matched_set = set()
    slide_frame_count: Dict[int, int] = {}
    confidence_counts: Dict[str, int] = {}
    unmatched_frame_count = 0

    for entry in slide_map:
        sb = entry.get("matched_slide", 0)
        conf = entry.get("match_confidence", "none")
        confidence_counts[conf] = confidence_counts.get(conf, 0) + 1

        if sb and sb > 0:
            matched_set.add(sb)
            slide_frame_count[sb] = slide_frame_count.get(sb, 0) + 1
        else:
            unmatched_frame_count += 1

    all_slides = set(range(1, total_slides + 1))
    unmatched_slides = all_slides - matched_set
    coverage_pct = (len(matched_set) / total_slides * 100) if total_slides > 0 else 0.0

    return {
        "total_slides": total_slides,
        "matched_slides": sorted(matched_set),
        "unmatched_slides": sorted(unmatched_slides),
        "coverage_pct": round(coverage_pct, 1),
        "unmatched_frames": unmatched_frame_count,
        "total_frames": len(slide_map),
        "slide_frame_count": dict(sorted(slide_frame_count.items())),
        "confidence_summary": confidence_counts,
    }


# ─────────────────────────────────────────────
# 음성·화면 파이프라인 병렬 실행 (속도 업그레이드)
# ─────────────────────────────────────────────

def run_pipeline_parallel(
    video_path: str,
    audio_path: str,
    api_key: str,
    *,
    check_audio: bool = True,
    check_screen: bool = True,
    # 화면 파이프라인 파라미터
    sample_rate: float = 1.5,
    diff_threshold: float = 25.0,
    screen_batch_size: int = 2,
    screen_max_workers: int = 4,
    # 음성 파이프라인 파라미터
    model: str = "gpt-5.4",
    domain_hint: str = "",
    context_window: int = 2,
    use_sentence_merge: bool = True,
    audio_batch_size: int = 40,
    audio_max_workers: int = 3,
    stt_chunk_workers: int = 3,
    storyboard_images: Optional[List[str]] = None,
    # ★ 캐싱 (안정성 업그레이드)
    stt_cache_dir: Optional[str] = None,
    start_sb_idx: int = 0,
    # ★ 정확도 v2 옵션
    narration_texts: Optional[List[str]] = None,
    verify_pass: bool = True,
    confidence_threshold: str = "medium",
    scene_change_threshold: float = 60.0,
    # ★ v3: 4항목 전면 검증용 슬라이드 메타데이터 (스토리보드-영상 QA 스킬 정합)
    slide_metadata: Optional[List[dict]] = None,
    # ★ v4: 안정화(애니메이션 종료) 캡처 모드
    stability_mode: bool = True,
    stability_motion_threshold: float = 2.0,
    stability_min_seconds: float = 1.2,
    stability_check_interval: float = 0.4,
    # ★ v10: 매칭 전용 중앙 크롭 이미지 (영상-스토리보드 1:1 매칭 정확도 향상)
    storyboard_match_images: Optional[List[str]] = None,
) -> Dict[str, Any]:
    """
    음성·화면 두 파이프라인을 **동시에** 실행합니다.

    ★ 안정성 옵션:
    - stt_cache_dir: 지정 시 Whisper 전사 결과를 파일 해시 기반으로 캐싱.
                     같은 영상·힌트 조합이면 즉시 재사용.

    ★ 정확도 v2 옵션:
    - narration_texts: 스토리보드 슬라이드별 내레이션 텍스트 리스트.
                       제공되면 STT 결과를 내레이션과 대조 검증 (단순 맞춤법 → 일치도 비교).
                       Whisper prompt에도 핵심 키워드를 주입하여 인식률 향상.
    - verify_pass: True면 화면 검사 1차 결과를 2차 검증하여 false positive 감소.
    - confidence_threshold: "high" | "medium" | "low". 내레이션 비교 시 이 수준 이상만 결과에 포함.
    - scene_change_threshold: 장면 전환 감지 임계값. 낮을수록 민감 (기본 60.0).

    Returns
    -------
    dict : {"audio": [...], "screen": [...], "slide_map": [...], "slide_coverage": {...},
            "errors": {...}, "stt_cached": bool, "last_sb_idx": int, "narration_mode": bool}
    """

    errors: Dict[str, str] = {}
    stt_cached = False
    use_narration = bool(narration_texts)

    def _screen_pipeline() -> Tuple[List[dict], List[dict], int]:
        if not check_screen:
            return [], [], start_sb_idx
        try:
            frames = extract_and_filter_frames(
                video_path, sample_rate, diff_threshold,
                scene_change_threshold=scene_change_threshold,
                stability_mode=stability_mode,
                stability_motion_threshold=stability_motion_threshold,
                stability_min_seconds=stability_min_seconds,
                stability_check_interval=stability_check_interval,
            )
            if not frames:
                return [], [], start_sb_idx
            return spell_check_frames(
                frames, api_key,
                batch_size=screen_batch_size,
                model=model,
                max_workers=screen_max_workers,
                storyboard_images=storyboard_images,
                storyboard_match_images=storyboard_match_images,   # ★ v10
                start_sb_idx=start_sb_idx,
                verify_pass=verify_pass,
                slide_metadata=slide_metadata,
            )
        except Exception as e:
            errors["screen"] = str(e)
            return [], [], start_sb_idx

    def _audio_pipeline() -> List[dict]:
        nonlocal stt_cached
        if not check_audio:
            return []
        try:
            segments = None

            # 1) 캐시 조회
            cache_key = None
            if stt_cache_dir:
                try:
                    cache_key = compute_stt_cache_key(video_path, domain_hint)
                    cached = load_stt_cache(stt_cache_dir, cache_key)
                    if cached:
                        segments = cached
                        stt_cached = True
                        print(f"⚡ STT 캐시 적중 — 전사 단계 건너뜀 ({len(cached)} segs)")
                except Exception as e:
                    print(f"STT 캐시 조회 실패 (무시하고 계속): {e}")

            # 2) 캐시 미스 → 정상 전사
            if segments is None:
                if not extract_audio(video_path, audio_path):
                    errors["audio"] = "오디오 추출 실패"
                    return []
                segments = transcribe_audio(
                    audio_path, api_key,
                    domain_hint=domain_hint,
                    max_workers=stt_chunk_workers,
                    narration_texts=narration_texts,  # ★ v2: Whisper prompt 강화
                )
                if stt_cache_dir and cache_key and segments:
                    save_stt_cache(stt_cache_dir, cache_key, segments)

            if not segments:
                return []

            # ★ 정확도 v2: 내레이션이 있으면 대조 비교, 없으면 기존 맞춤법 검사
            if use_narration:
                results = narration_compare_segments(
                    segments, api_key,
                    narration_texts=narration_texts,
                    context_window=context_window,
                    model=model,
                    use_sentence_merge=use_sentence_merge,
                    batch_size=audio_batch_size,
                    max_workers=audio_max_workers,
                    confidence_threshold=confidence_threshold,
                )
            else:
                results = spell_check_segments(
                    segments, api_key,
                    context_window=context_window,
                    model=model,
                    use_sentence_merge=use_sentence_merge,
                    batch_size=audio_batch_size,
                    max_workers=audio_max_workers,
                )
            return attach_frames_to_audio_results(results, video_path)
        except Exception as e:
            errors["audio"] = str(e)
            return []

    # 두 파이프라인 병렬 실행
    with ThreadPoolExecutor(max_workers=2) as ex:
        fut_screen = ex.submit(_screen_pipeline)
        fut_audio  = ex.submit(_audio_pipeline)
        screen_results, slide_map, last_sb_idx = fut_screen.result()
        audio_results  = fut_audio.result()

    # ★ 정확도 v2: 슬라이드 커버리지 리포트 생성
    slide_coverage = {}
    if storyboard_images and slide_map:
        slide_coverage = compute_slide_coverage(slide_map, len(storyboard_images))

    return {
        "audio":  audio_results,
        "screen": screen_results,
        "slide_map": slide_map,
        "slide_coverage": slide_coverage,
        "errors": errors,
        "stt_cached": stt_cached,
        "last_sb_idx": last_sb_idx,
        "narration_mode": use_narration,
    }


# ─────────────────────────────────────────────
# 내부 헬퍼
# ─────────────────────────────────────────────

def _strip_json_fences(text):
    """마크다운 코드 펜스(```json ... ```)를 제거합니다."""
    if text.startswith("```json"):
        text = text[7:]
    elif text.startswith("```"):
        text = text[3:]
    if text.endswith("```"):
        text = text[:-3]
    return text.strip()


def _red_to_html(text):
    """<red>단어</red> 태그를 인라인 HTML 강조 스타일로 변환합니다."""
    return (text
            .replace("<red>", "<span style='color:red; font-weight:bold;'>")
            .replace("</red>", "</span>"))
